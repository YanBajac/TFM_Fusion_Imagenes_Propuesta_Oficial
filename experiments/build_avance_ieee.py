# -*- coding: utf-8 -*-
"""Reconstruye la presentación de avance con estilo minimalista tipo paper/IEEE."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from PIL import Image

INK=RGBColor(0x1A,0x1A,0x1A); GRAY=RGBColor(0x55,0x55,0x55)
LIGHT=RGBColor(0x8A,0x8A,0x8A); ACC=RGBColor(0x2F,0x5B,0x7C)
RULE=RGBColor(0xCF,0xCF,0xCF)
FONT="Arial"
SW, SH = Inches(13.333), Inches(7.5)
MX=Inches(0.85)
FIGDIR="docs/figures/"

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]
TOTAL=21

def slide():
    return prs.slides.add_slide(BLANK)

def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    b=s.shapes.add_textbox(x,y,w,h); tf=b.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.paragraphs[0].alignment=align
    return b,tf

def run(p, text, size, color=INK, bold=False, italic=False, font=FONT, spacing=None):
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return r

def line(s, x1, y1, x2, y2, color=RULE, w=0.75):
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1,y1,x2,y2)
    ln.line.color.rgb=color; ln.line.width=Pt(w); return ln

def footer(s, page, left="Método óptimo multiescala · Fusión VIS+IR — UCOM"):
    line(s, MX, Inches(7.02), SW-MX, Inches(7.02), RULE, 0.75)
    b,tf=tb(s, MX, Inches(7.08), Inches(8), Inches(0.3))
    run(tf.paragraphs[0], left, 9, LIGHT)
    b2,tf2=tb(s, SW-MX-Inches(2.5), Inches(7.08), Inches(2.5), Inches(0.3), align=PP_ALIGN.RIGHT)
    run(tf2.paragraphs[0], f"{page} / {TOTAL}", 9, LIGHT)

def header(s, tag, title, page):
    b,tf=tb(s, MX, Inches(0.52), SW-2*MX, Inches(0.35))
    run(tf.paragraphs[0], tag.upper(), 11, ACC, bold=True)
    b2,tf2=tb(s, MX, Inches(0.86), SW-2*MX, Inches(0.7))
    run(tf2.paragraphs[0], title, 26, INK, bold=True)
    footer(s, page)

def subhead(tf, text, first=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before=Pt(0 if first else 10); p.space_after=Pt(3)
    run(p, text, 13, ACC, bold=True)
    return p

def bullet(tf, text, size=14, first=False, color=INK, dash="—  "):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after=Pt(5); p.line_spacing=1.05
    run(p, dash, size, ACC, bold=True)
    run(p, text, size, color)
    return p

def para(tf, text, size=14, first=False, color=INK, bold=False, align=PP_ALIGN.LEFT):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(6); p.line_spacing=1.1
    run(p, text, size, color, bold=bold)
    return p

def figure(s, path, x, y, w_max, h_max, caption=None):
    iw,ih=Image.open(path).size; ar=iw/ih
    w=w_max; h=int(w/ar)
    if h>h_max: h=h_max; w=int(h*ar)
    left=int(x+(w_max-w)/2)
    s.shapes.add_picture(path, left, y, width=w, height=h)
    if caption:
        b,tf=tb(s, x, y+h+Inches(0.08), w_max, Inches(0.5), align=PP_ALIGN.CENTER)
        run(tf.paragraphs[0], caption, 10.5, LIGHT, italic=True)
    return h

# ============================ S1 — Portada ============================
s=slide()
line(s, MX, Inches(1.7), SW-MX, Inches(1.7), INK, 1.2)
b,tf=tb(s, MX, Inches(1.0), SW-2*MX, Inches(0.5))
run(tf.paragraphs[0], "UNIVERSIDAD COMUNERA — UCOM   ·   Maestría en Ciencias de Datos", 13, GRAY, bold=True)
b,tf=tb(s, MX, Inches(2.2), SW-2*MX, Inches(1.6))
run(tf.paragraphs[0], "Fusión de imágenes infrarrojas y visibles", 38, INK, bold=True)
p=tf.add_paragraph(); run(p, "mediante morfología matemática", 38, INK, bold=True)
b,tf=tb(s, MX, Inches(3.9), SW-2*MX, Inches(0.5))
run(tf.paragraphs[0], "Un método óptimo multiescala con elementos estructurantes circulares y lineales, optimizado por PSO", 18, ACC)
line(s, MX, Inches(5.4), SW-MX, Inches(5.4), RULE, 0.75)
b,tf=tb(s, MX, Inches(5.6), SW-2*MX, Inches(1.1))
run(tf.paragraphs[0], "Autores:  Lic. Juan Pablo Bazán   ·   Ing. Yan Bajac", 13, INK)
p=tf.add_paragraph(); run(p, "Director:  D.Sc. Julio César Mello", 13, INK)
p=tf.add_paragraph(); p.space_before=Pt(6); run(p, "Avance de tesis   ·   Asunción, 2026", 11, LIGHT)

# ============================ S2 — Agenda ============================
s=slide(); header(s, "Contenido", "Agenda", 2)
items=["Problema y objetivos","Marco teórico — método multiescala","Diseño experimental",
 "Hallazgo metodológico (refinamiento)","Resultados cuantitativos","Análisis estadístico",
 "Conclusiones y próximos pasos"]
b,tf=tb(s, MX, Inches(1.7), SW-2*MX, Inches(5.0))
for i,it in enumerate(items):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after=Pt(11)
    run(p, f"{i+1:02d}   ", 16, ACC, bold=True)
    run(p, it, 16, INK)

# ============================ S3 — Problema ============================
s=slide(); header(s,"1 · Problema","Problema y contexto",3)
b,tf=tb(s, MX, Inches(1.7), Inches(11.6), Inches(5.0))
subhead(tf,"Por qué fusionar VIS + IR", first=True)
bullet(tf,"VIS: alta resolución textural y color, sensible a iluminación adversa (noche, niebla, humo).")
bullet(tf,"IR: robusto a oscuridad y oclusores atmosféricos, pero baja resolución y poca textura.")
bullet(tf,"La fusión integra ambas modalidades en una representación útil para vigilancia, conducción autónoma e inspección industrial.")
subhead(tf,"Brecha en la literatura")
bullet(tf,"La morfología matemática preserva geometría con bajo costo computacional, pero los trabajos previos suelen restringirse a una única escala o un único elemento estructurante.")
bullet(tf,"No existe un estudio sistemático que articule SE × niveles × radio, contrastado contra baselines clásicos.")
para(tf,"Esta tesis aborda esa brecha con un benchmark factorial de 36 configuraciones × 20 pares VIS/IR.", size=14, bold=True)

# ============================ S4 — Objetivos ============================
s=slide(); header(s,"1 · Objetivos","Objetivos",4)
b,tf=tb(s, MX, Inches(1.7), Inches(11.6), Inches(1.4))
subhead(tf,"Objetivo general", first=True)
para(tf,"Proponer y evaluar un método óptimo multiescala de fusión VIS+IR (Top-Hat con elementos estructurantes circulares y lineales, optimizado por PSO), comparándolo con el método anterior (Torre Top-Hat) y con baselines clásicos, y analizando su impacto en una tarea de detección, sobre el TNO Image Fusion Dataset.")
b,tf=tb(s, MX, Inches(3.5), Inches(11.6), Inches(3.2))
subhead(tf,"Objetivos específicos", first=True)
esp=["Formular el método óptimo multiescala: banco de 5 SE (disco + 4 líneas) y cascada de n escalas.",
 "Optimizar automáticamente sus hiperparámetros (n, r, m) mediante enjambre de partículas (PSO).",
 "Comparar con el método anterior (Torre Top-Hat) y los baselines con 12 métricas y pruebas no paramétricas.",
 "Analizar el impacto de la fusión en una tarea de detección de objetos (YOLO)."]
for i,e in enumerate(esp):
    p=tf.add_paragraph(); p.space_after=Pt(7)
    run(p, f"{i+1:02d}   ", 14, ACC, bold=True); run(p, e, 14, INK)

# ============================ S5 — Marco teórico ============================
s=slide(); header(s,"2 · Marco teórico","Transformada Top-Hat",5)
b,tf=tb(s, MX, Inches(1.7), Inches(11.6), Inches(5.0))
subhead(tf,"White Top-Hat (WTH)", first=True)
bullet(tf,"WTH(f, b) = f − γ(f, b),  con γ = apertura morfológica = δ(ε(f, b), b).")
bullet(tf,"Aísla estructuras brillantes locales menores que el elemento estructurante b.")
bullet(tf,"Es siempre no-negativa; se anula donde la imagen es plana a la escala del SE.")
subhead(tf,"Decisión de diseño central: el elemento estructurante (SE)")
bullet(tf,"La geometría (disco, cuadrado, cruz) y el radio del SE definen qué orientaciones, simetrías y escalas se preservan o eliminan.")
bullet(tf,"La Torre Top-Hat encadena varias de estas transformadas con SE de radio creciente para extraer detalles a múltiples escalas.")

# ============================ S6 — Método (figura) ============================
s=slide(); header(s,"2 · Método propuesto","Método óptimo multiescala",6)
figure(s, "docs/figures/disenos/C_esquema_vectorial.png", MX, Inches(1.85), Inches(7.4), Inches(3.0))
b,tf=tb(s, Inches(8.5), Inches(1.7), Inches(3.95), Inches(4.9))
subhead(tf,"Operador", first=True)
bullet(tf,"Banco de 5 SE por escala: 1 disco + 4 líneas (0/45/90/135°), combinados por máximo.", size=12.5)
bullet(tf,"Cascada de n escalas (radios r·i); agregación por máximo entre escalas.", size=12.5)
subhead(tf,"Regla de fusión")
p=tf.add_paragraph(); p.space_after=Pt(4)
run(p,"F = I_base + m·(WTH_M+WTHS_M) − m·(BTH_M+BTHS_M)", 12, ACC, bold=True)
subhead(tf,"Optimización (PSO)")
para(tf,"Los hiperparámetros (n, r, m) se ajustan por enjambre de partículas → n=6, r≈2,89, m=0,10. La Torre Top-Hat (un disco por escala) es el método anterior de comparación.", size=12, color=GRAY)

# ============================ S7 — Diseño (stats) ============================
s=slide(); header(s,"3 · Diseño experimental","Diseño experimental",7)
stats=[("20","pares VIS/IR"),("36","configuraciones Top-Hat"),("11","métodos comparados"),("12","métricas no-reference")]
cw=Inches(2.7); gap=Inches(0.45); x0=MX; y0=Inches(1.8)
for i,(num,lab) in enumerate(stats):
    x=x0+i*(cw+gap)
    b,tf=tb(s, x, y0, cw, Inches(1.2), align=PP_ALIGN.LEFT)
    run(tf.paragraphs[0], num, 44, ACC, bold=True)
    p=tf.add_paragraph(); run(p, lab, 12, GRAY)
    if i<3: line(s, x+cw+Inches(0.12), y0+Inches(0.15), x+cw+Inches(0.12), y0+Inches(1.15), RULE, 0.75)
b,tf=tb(s, MX, Inches(3.6), Inches(11.6), Inches(3.0))
subhead(tf,"Estudio 1 — Benchmark factorial de la Torre Top-Hat", first=True)
bullet(tf,"SE ∈ {disco, cuadrado, cruz} × Niveles ∈ {2,3,4,5} × r₀ ∈ {2,3,5}.  Total: 36 × 20 = 720 fusiones.")
subhead(tf,"Estudio 2 — Comparación inter-método")
bullet(tf,"Promedio · Pirámide Laplace · Curvelet · 4 Top-Hat WTH · 4 WTH+BTH, sobre 12 métricas (6 clásicas + 6 estándar de calidad).")

# ============================ S8 — Dataset (figura) ============================
s=slide(); header(s,"3 · Dataset","TNO Image Fusion Dataset",8)
b,tf=tb(s, MX, Inches(1.65), Inches(11.6), Inches(0.9))
para(tf,"TNO Multiband Image Data Collection (Toet, 2014). 20 pares seleccionados: vehículos blindados, soldados tras humo, escenas urbanas a contraluz, vegetación e interiores industriales.", first=True, size=14)
figure(s, FIGDIR+"fig_qualitative_comparison.png", MX, Inches(2.7), Inches(11.6), Inches(3.5),
       caption="Ejemplos de pares VIS / IR y comparación cualitativa de los principales métodos.")

# ============================ S9 — Refinamiento ident ============================
s=slide(); header(s,"4 · Refinamiento metodológico","Identificación del sesgo",9)
b,tf=tb(s, MX, Inches(1.6), Inches(11.6), Inches(0.7))
para(tf,"Durante la revisión iterativa del código se identificó un sesgo metodológico en la regla de fusión.", first=True, size=14, bold=True)
b,tf=tb(s, MX, Inches(2.5), Inches(5.6), Inches(4.0))
subhead(tf,"Versión inicial (sesgada)", first=True)
bullet(tf,"Max-actividad aplicada a TODAS las capas, incluida la base.", size=12.5)
bullet(tf,"La base se elegía píxel a píxel entre VIS e IR → discontinuidades de iluminación.", size=12.5)
bullet(tf,"La información mutua (MI) quedaba inflada artificialmente.", size=12.5)
b,tf=tb(s, Inches(7.0), Inches(2.5), Inches(5.45), Inches(4.0))
subhead(tf,"Versión corregida", first=True)
bullet(tf,"Reglas distintas: detalle por max-actividad, base por PROMEDIO.", size=12.5)
bullet(tf,"Coherente con Burt-Adelson y la fusión multiescala clásica.", size=12.5)
bullet(tf,"Verificación: la descomposición sin BTH reconstruye f exactamente (error 0.0).", size=12.5)
b,tf=tb(s, MX, Inches(6.2), Inches(11.6), Inches(0.6))
para(tf,"Este refinamiento cambió las conclusiones cuantitativas de la tesis.", first=True, size=12.5, color=GRAY, italic_bold=False) if False else para(tf,"Este refinamiento cambió las conclusiones cuantitativas de la tesis.", first=True, size=12.5, color=GRAY)

# ============================ S10 — Refinamiento impacto ============================
s=slide(); header(s,"4 · Refinamiento metodológico","Impacto en las conclusiones",10)
b,tf=tb(s, MX, Inches(1.9), Inches(11.6), Inches(4.0))
subhead(tf,"Lectura honesta", first=True)
para(tf,"La Torre Top-Hat ya no domina la información mutua por construcción del sesgo. La Pirámide de Laplace pasa a liderar 5 de 6 métricas clásicas en valor medio (Curvelet gana MG).")
para(tf,"El promedio simple, por su naturaleza lineal F=(V+I)/2, conserva información de ambas fuentes y compite en MI per-imagen. La Torre Top-Hat queda en una posición media-alta como método clásico interpretable y competitivo.")

# ============================ S11 — Resultados promedios (tabla) ============================
s=slide(); header(s,"5 · Resultados","Promedios por método",11)
figure(s, FIGDIR+"_pres/tabla_promedios.png", MX, Inches(1.7), Inches(11.6), Inches(3.2))
b,tf=tb(s, MX, Inches(5.2), Inches(11.6), Inches(1.6))
bullet(tf,"La Pirámide de Laplace lidera 5 de 6 métricas clásicas; Top-Hat disco L5 encabeza entre las configuraciones morfológicas (EN, SD, FE, MG).", first=True, size=12.5)
bullet(tf,"Las diferencias absolutas entre las 4 configuraciones Top-Hat son modestas (< 0,06 en EN).", size=12.5)

# ============================ S12 — Distribución (figura) ============================
s=slide(); header(s,"5 · Resultados","Distribución por método",12)
figure(s, FIGDIR+"fig_boxplot_metrics.png", MX, Inches(1.65), Inches(11.6), Inches(4.0),
       caption="Diagramas de caja por método (n = 20 pares VIS/IR).")
b,tf=tb(s, MX, Inches(6.15), Inches(11.6), Inches(0.7))
para(tf,"Las diferencias en EN, SD, FE y MG entre métodos son sólidas; en MI la dispersión es alta y favorece al Promedio.", first=True, size=12, color=GRAY)

# ============================ S13 — Friedman + Ranking ============================
s=slide(); header(s,"6 · Análisis estadístico","Friedman y ranking global",13)
b,tf=tb(s, MX, Inches(1.65), Inches(5.7), Inches(0.4))
run(tf.paragraphs[0],"Prueba de Friedman global", 13, ACC, bold=True)
figure(s, FIGDIR+"_pres/tabla_friedman.png", MX, Inches(2.1), Inches(5.7), Inches(2.7))
b,tf=tb(s, MX, Inches(4.9), Inches(5.7), Inches(0.8))
para(tf,"Las doce métricas resultan altamente significativas (p << 0,001); se muestran seis.", first=True, size=11.5, color=GRAY)
b,tf=tb(s, Inches(6.9), Inches(1.65), Inches(5.5), Inches(0.4))
run(tf.paragraphs[0],"Ranking global (menor es mejor)", 13, ACC, bold=True)
figure(s, FIGDIR+"fig_ranking_global_12metricas.png", Inches(6.9), Inches(2.1), Inches(5.55), Inches(3.6))
b,tf=tb(s, Inches(6.9), Inches(5.85), Inches(5.5), Inches(0.8))
para(tf,"Pirámide de Laplace lidera el agregado (4,42); la mejor Top-Hat WTH es disco L5 (5,00), muy próxima.", first=True, size=11.5, color=GRAY)

# ============================ S14 — Benchmark heatmap ============================
s=slide(); header(s,"6 · Análisis estadístico","Benchmark exhaustivo (36 configuraciones)",14)
figure(s, FIGDIR+"fig_heatmap_tophat.png", MX, Inches(1.65), Inches(11.6), Inches(4.0),
       caption="Mapas de calor de las métricas sobre las 36 configuraciones de la Torre Top-Hat.")
b,tf=tb(s, MX, Inches(6.15), Inches(11.6), Inches(0.7))
para(tf,"EN/SD/FE crecen con niveles y radio; MG óptimo en disco L5 r3; MI favorece configuraciones someras. Diferencias internas modestas (robustez de hiperparámetros).", first=True, size=12, color=GRAY)

# ============================ S15 — Métricas estándar ============================
s=slide(); header(s,"5 · Resultados","Método óptimo multiescala vs. baselines",15)
figure(s, FIGDIR+"fig_metodo_optimo_multiescala.png", MX, Inches(1.65), Inches(11.6), Inches(4.0),
       caption="Método óptimo multiescala frente al método anterior y a los baselines en las métricas estándar de calidad.")
b,tf=tb(s, MX, Inches(6.15), Inches(11.6), Inches(0.7))
para(tf,"Lidera Qabf (0,514) y SCD (1,453), menor Nabf entre métodos no triviales (0,065) y 2º en SSIM (0,775); cede en contraste (SD, EN). Óptimo PSO: n=6, r≈2,89, m=0,10.", first=True, size=12, color=GRAY)

# ============================ S16 — Efecto BTH ============================
s=slide(); header(s,"5 · Resultados","Efecto de la variante Black Top-Hat",16)
figure(s, FIGDIR+"fig_efecto_black_top_hat.png", MX, Inches(1.65), Inches(11.6), Inches(3.9),
       caption="Cambio relativo al añadir BTH (verde = mejora; rojo = empeora).")
b,tf=tb(s, MX, Inches(6.05), Inches(11.6), Inches(0.8))
para(tf,"Añadir Black Top-Hat sube contraste y actividad (EN, SD, SF, MG) pero degrada la calidad (Qabf, SSIM, VIF, MI) y multiplica los artefactos (Nabf ×5). No se recomienda por defecto.", first=True, size=12, color=GRAY)

# ============================ S17 — Detección YOLO ============================
s=slide(); header(s,"6 · Evaluación orientada a tarea","Detección con YOLO",17)
figure(s, FIGDIR+"fig_deteccion_yolo.png", MX, Inches(1.65), Inches(11.6), Inches(3.9),
       caption="Detecciones de YOLOv8n (inferencia) por modalidad y por método.")
b,tf=tb(s, MX, Inches(6.05), Inches(11.6), Inches(0.85))
para(tf,"En VIS se detecta solo 1 persona; con fusión sube a 13 (WTH) y 16 (+BTH). VIS favorece vehículos, IR/fusión personas; WTH equilibra ambas. Diferencias entre métodos no significativas (n=20); mAP pendiente con dataset etiquetado.", first=True, size=12, color=GRAY)

# ============================ S18 — Discusión ============================
s=slide(); header(s,"7 · Discusión","Discusión integrada",18)
disc=[("Método óptimo multiescala: mayor calidad de fusión","Lidera Qabf (0,514) y SCD (1,453), menor Nabf entre métodos no triviales (0,065) y 2º en SSIM (0,775). Supera al método anterior en bordes, correlación y artefactos."),
 ("Pirámide de Laplace y Torre Top-Hat","Laplace encabeza el ranking agregado de 12 métricas (4,42) por contraste (SD) y VIF; la Torre Top-Hat (anterior, disco L5) es 2ª (5,00)."),
 ("Sesgo intrínseco de la métrica MI","El promedio gana MI per-imagen por su naturaleza lineal. Conviene reportar MI siempre acompañada de métricas de calidad global."),
 ("No hay método universalmente óptimo","La elección depende del criterio operativo: equilibrio global (Laplace), fidelidad estructural / interpretabilidad (Top-Hat) o trazabilidad lineal (Promedio).")]
b,tf=tb(s, MX, Inches(1.7), Inches(11.6), Inches(5.0))
for i,(h,t) in enumerate(disc):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_before=Pt(0 if i==0 else 9); p.space_after=Pt(2)
    run(p, f"{i+1:02d}   ", 14, ACC, bold=True); run(p, h, 14, INK, bold=True)
    p2=tf.add_paragraph(); p2.space_after=Pt(2); run(p2, t, 12.5, GRAY)

# ============================ S19 — Conclusiones ============================
s=slide(); header(s,"Cierre","Conclusiones",19)
conc=["El método óptimo multiescala propuesto logra la mayor calidad de fusión: lidera Qabf (0,514), SCD (1,453) y Nabf (0,065), y es 2º en SSIM (0,775); valida los SE lineales, la cascada y el ajuste por PSO.",
 "Supera al método anterior (Torre Top-Hat, modo WTH) en bordes, correlación y artefactos; la Torre, a su vez, supera a Laplace en fidelidad estructural (SSIM, SCD).",
 "La Pirámide de Laplace logra el mejor ranking agregado (4,42) por contraste y VIF; Curvelet es el más débil al penalizar artefactos (Nabf).",
 "Diferencias estadísticamente robustas: Friedman p << 0,001 en las 12 métricas; Wilcoxon con corrección de Holm y tamaño de efecto.",
 "La variante Black Top-Hat aumenta el contraste pero degrada la calidad y multiplica los artefactos (Nabf ×5): no se recomienda por defecto."]
b,tf=tb(s, MX, Inches(1.7), Inches(11.6), Inches(5.0))
for i,c in enumerate(conc):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after=Pt(9); p.line_spacing=1.05
    run(p, f"{i+1:02d}   ", 14, ACC, bold=True); run(p, c, 14, INK)

# ============================ S20 — Próximos pasos ============================
s=slide(); header(s,"Cierre","Próximos pasos",20)
steps=[("Corto plazo","Probar reglas que usen el Black Top-Hat de forma selectiva sin introducir artefactos; cerrar la evaluación de detección con dataset etiquetado."),
 ("Mediano plazo","Reglas de fusión adaptativas (saliencia, regiones); comparación contra modelos profundos (FusionGAN, U2Fusion) como referencia."),
 ("Antes de la defensa","Formato definitivo según normas UCOM; revisión integral con el director; preparación de la defensa.")]
b,tf=tb(s, MX, Inches(1.8), Inches(11.6), Inches(4.5))
for i,(h,t) in enumerate(steps):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_before=Pt(0 if i==0 else 12); p.space_after=Pt(3)
    run(p, h, 14, ACC, bold=True)
    p2=tf.add_paragraph(); p2.space_after=Pt(2); run(p2, t, 13.5, INK)

# ============================ S21 — Gracias ============================
s=slide()
line(s, MX, Inches(2.9), SW-MX, Inches(2.9), INK, 1.2)
b,tf=tb(s, MX, Inches(3.1), SW-2*MX, Inches(1.0))
run(tf.paragraphs[0], "Gracias", 40, INK, bold=True)
b,tf=tb(s, MX, Inches(4.2), SW-2*MX, Inches(0.5))
run(tf.paragraphs[0], "Discusión y comentarios", 16, ACC)
b,tf=tb(s, MX, Inches(5.2), SW-2*MX, Inches(1.0))
run(tf.paragraphs[0], "Lic. Juan Pablo Bazán   ·   Ing. Yan Bajac", 12, INK)
p=tf.add_paragraph(); run(p, "Director: D.Sc. Julio César Mello", 12, INK)
p=tf.add_paragraph(); p.space_before=Pt(4); run(p, "UCOM — Maestría en Ciencias de Datos · 2026", 10.5, LIGHT)

prs.save("docs/Tesis_Avance_Presentacion.pptx")
print("OK slides:", len(prs.slides._sldIdLst))
