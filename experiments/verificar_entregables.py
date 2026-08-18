# -*- coding: utf-8 -*-
"""Verificador unico de los cuatro entregables contra los CSV.

Motivo: durante el proyecto se repitio el mismo patron —texto y figuras copiados
entre documentos que despues divergen de los datos—. Aparecio en el libro, en el
deck, en los montajes cualitativos y en el README. Este script cierra ese agujero:
un solo comando que falla si cualquiera de los cuatro contradice a los CSV.

Todos los valores esperados se DERIVAN de experiments/results/metrics_reports/.
Nada esta escrito a mano, salvo la lista de afirmaciones retiradas (§4), que por
definicion no puede derivarse de los datos.

Uso:   .venv\Scripts\python.exe -X utf8 experiments/verificar_entregables.py
Salida: informe por consola; codigo de salida 1 si hay algun fallo.
"""
import ast
import hashlib
import re
import sys
import unicodedata
import zipfile
from pathlib import Path

import fitz
import pandas as pd

RAIZ = Path(__file__).resolve().parent.parent
REP = RAIZ / 'experiments' / 'results' / 'metrics_reports'
DOCS = RAIZ / 'docs'

fallos = []
avisos = []


def ok(cond, msg, blando=False):
    if cond:
        print(f'  OK    {msg}')
    elif blando:
        print(f'  AVISO {msg}')
        avisos.append(msg)
    else:
        print(f'  FALLA {msg}')
        fallos.append(msg)
    return bool(cond)


def coma(v, nd):
    return f'{v:.{nd}f}'.replace('.', ',')


def punto(v, nd):
    return f'{v:.{nd}f}'


def plano(s):
    """Normaliza espacios: la extraccion de PDF corta frases y numeros en saltos."""
    return re.sub(r'\s+', ' ', s)


def sin_tildes(s):
    s = unicodedata.normalize('NFKD', s)
    return ''.join(c for c in s if not unicodedata.combining(c))


# ----------------------------------------------------------------- documentos
def texto_pdf(p):
    with fitz.open(str(p)) as d:
        return plano('\n'.join(pg.get_text() for pg in d)), d.page_count


DOCUMENTOS = {}
for clave, ruta in (('libro', DOCS / 'Tesis_Borrador_V3.pdf'),
                    ('deck', DOCS / 'Tesis_Defensa_Presentacion.pdf'),
                    ('avances', DOCS / 'Avances_Tesis.pdf')):
    if ruta.exists():
        t, n = texto_pdf(ruta)
        DOCUMENTOS[clave] = {'texto': t, 'paginas': n, 'ruta': ruta}
readme = RAIZ / 'README.md'
if readme.exists():
    DOCUMENTOS['readme'] = {'texto': plano(readme.read_text(encoding='utf-8')),
                            'paginas': None, 'ruta': readme}

print('=== documentos ===')
for k, v in DOCUMENTOS.items():
    pg = f'{v["paginas"]} pag.' if v['paginas'] else f'{len(v["texto"])} car.'
    print(f'  {k:8} {pg:12} {v["ruta"].name}')
faltan = [k for k in ('libro', 'deck', 'avances', 'readme') if k not in DOCUMENTOS]
ok(not faltan, f'los cuatro entregables existen (faltan: {faltan or "ninguno"})')


def contiene(clave, *variantes):
    """True si alguna variante aparece en el documento COMO CIFRA ENTERA, no como subcadena.

    La version anterior hacia `v in t` a secas, y eso da falsos positivos con las variantes de dos
    decimales, que son cortas. Paso de verdad: al publicar la tabla de diez columnas de la lamina 12,
    que trae MI_vis de la Ratio Pyramid = 0,949, el chequeo de los mAP de LLVIP dejo de reclamar
    Curvelet (0,9403) y DWT (0,9394), porque los dos buscaban «0,94» y «0,949» lo contiene. El aviso
    paso de tres metodos faltantes a uno y parecia una mejora: no lo era, era el chequeo aflojandose.

    Ahora se exige que ni antes ni despues del numero haya otro digito, ni un separador decimal
    seguido de digito. Asi «0,94» ya no se satisface con «0,949» ni con «10,94», y sigue encontrandose
    en «0,94 » o «(0,94)».
    """
    t = DOCUMENTOS[clave]['texto']
    for v in variantes:
        for m in re.finditer(re.escape(v), t):
            antes, desp = t[:m.start()], t[m.end():]
            if antes and antes[-1].isdigit():
                continue                                 # es la cola de un numero mas largo
            if len(antes) >= 2 and antes[-1] in ',.' and antes[-2].isdigit():
                continue                                 # p. ej. buscar «94» dentro de «0,94»
            if desp and desp[0].isdigit():
                continue                                 # tiene mas decimales
            if len(desp) >= 2 and desp[0] in ',.' and desp[1].isdigit():
                continue                                 # es la parte entera de otro numero
            return True
    return False


# --------------------------------------------- 1. medias del benchmark
print('\n=== 1. medias por metodo (9 metricas clasicas) ===')
dm = pd.read_csv(REP / 'descriptive_means.csv').set_index('method')
M9 = ['EN', 'SD', 'FE', 'MG', 'MI_vis', 'MI_ir', 'SF', 'SSIM', 'PSNR']
for doc in ('libro', 'avances'):
    if doc not in DOCUMENTOS:
        continue
    falt = []
    for met in dm.index:
        for m in M9:
            v = float(dm.loc[met, m])
            # se acepta cualquier redondeo entre 2 y 4 decimales, con coma o punto
            if not contiene(doc, *[f(v, nd) for nd in (4, 3, 2) for f in (coma, punto)]):
                falt.append(f'{met}/{m}')
    ok(not falt, f'{doc}: las {len(dm)*9} medias aparecen'
                 + (f' — faltan {falt[:5]}' if falt else ''),
       blando=(doc == 'avances'))

# --------------------------------------------- 2. ranking
print('\n=== 2. ranking de rangos medios ===')
rk = pd.read_csv(REP / 'ranking_methods.csv', index_col=0)['avg_rank']
primero = rk.idxmin()
ok(primero == 'Propuesta_Novedosa',
   f'el CSV pone primera a la propuesta ({primero}, {rk.min():.3f})')
for doc in DOCUMENTOS:
    falt = [f'{i}={v:.2f}' for i, v in rk.items()
            if not contiene(doc, coma(v, 2), punto(v, 2), coma(v, 3), punto(v, 3))]
    ok(not falt, f'{doc}: los {len(rk)} rangos globales aparecen'
                 + (f' — faltan {falt}' if falt else ''),
       blando=(doc in ('deck', 'readme')))
# la propuesta debe declararse primera, no segunda
for doc in DOCUMENTOS:
    malas = [s for s in ('segundo lugar del ranking', 'queda segunda, a 0,23',
                         'ocupa el segundo lugar (3,67)')
             if s in DOCUMENTOS[doc]['texto']]
    ok(not malas, f'{doc}: no declara a la propuesta en segundo lugar {malas}')

# --------------------------------------------- 3. deteccion
print('\n=== 3. deteccion ===')
# LLVIP se compara contra la MEDIA DE LAS CINCO SEMILLAS: es lo que los tres entregables publican
# desde que el estudio de repeticiones existe. La corrida de una semilla se conserva aparte, para
# los chequeos que hablan expresamente de ella.
ll_1sem = pd.read_csv(REP / 'detection_llvip_map.csv').set_index('method')
_sem_res = REP / 'semillas_llvip_resumen.csv'
ll = (pd.read_csv(_sem_res, index_col=0).rename(columns={'mAP50_media': 'mAP50'})
      if _sem_res.exists() else ll_1sem)
m3 = pd.read_csv(REP / 'detection_m3fd_map.csv').set_index('method')
com = pd.read_csv(REP / 'complementariedad_resumen.csv').set_index('entrada')
for doc in ('libro', 'deck', 'avances'):
    if doc not in DOCUMENTOS:
        continue
    falt = [f'LLVIP {i}' for i, v in ll['mAP50'].items()
            if not contiene(doc, coma(v, 3), punto(v, 3), coma(v, 2), punto(v, 2))]
    ok(not falt, f'{doc}: los mAP de LLVIP aparecen' + (f' — faltan {falt}' if falt else ''),
       blando=(doc == 'deck'))
# el par complementario: la mejor entrada es Ratio Pyramid, no la propuesta
par = ((m3['AP50_People'] + m3['AP50_Lamp']) / 2).sort_values(ascending=False)
ok(par.index[0] == 'RatioPiramide',
   f'el CSV pone a RatioPiramide primera en el par complementario ({par.iloc[0]:.3f})')
pos_prop = list(par.index).index('Propuesta_Novedosa') + 1
print(f'        la propuesta queda {pos_prop}.a de {len(par)} en el par complementario')
# el conteo por escena
p_prop = float(com.loc['Propuesta_Novedosa', 'pct_ambas'])
p_vis = float(com.loc['VIS', 'pct_ambas'])
ok(p_prop < p_vis, f'el CSV deja la propuesta por debajo del visible en el conteo '
                   f'por escena ({p_prop:.1f} % vs {p_vis:.1f} %)')

# --------------------------------------------- 4. afirmaciones retiradas
print('\n=== 4. afirmaciones retiradas (no pueden reaparecer) ===')
# Varias de estas frases SI aparecen legitimamente, negadas: el texto corregido dice
# «no hay patron espejo», «el infrarrojo no es ciego a las luces», «r = 1 no
# desactiva el banco». Por eso no basta buscar la cadena: hay que mirar si viene
# precedida de una negacion en la misma oracion.
RETIRADAS = [
    'la propuesta es la mejor fusión del estudio',
    'ciego a las luces', 'ciego a Lamp', 'el VIS es el espejo',
    'Complementariedad extrema', 'patrón espejo',
    'maximiza las nueve métricas', 'La configuración hallada por PSO',
    'desactiva el banco', 'contenido de bordes', 'entropía de bordes',
    'penaliza los artefactos', 'pocos artefactos',
    '36 configuraciones', '90 contrastes', '[0,05; 1,20]', '0,05–1,20',
    'seis personas frente a dos', 'cuatro luces que el IR',
    'Toet, A. (2014)', 'Mukhopadhyay y Chanda (2001)',
    # El mAP del visible y del infrarrojo en LLVIP estaban escritos a mano en la lectura
    # de la Tabla 9 del informe de avances, con valores de una corrida anterior, mientras
    # la tabla de arriba —generada desde el CSV— imprimia los vigentes: 0,813 y 0,971.
    # Se buscan como frase y no como cifra suelta porque 0,808 y 0,957 colisionan con
    # valores por imagen que si son vigentes (SSIM, Q0, QW, FE de varios CSV).
    'de 0,808 a la banda', 'más fuerte (0,957)',
    # y el argumento que los otros tres documentos ya retiraron
    'activa el banco completo',
    # El deck contaba «r = 1 en 16 de las 25 configuraciones» y «r = 25 ... aparece en 8». Es el
    # reparto del barrido publicado, que corre UNA semilla por celda; con 20 repeticiones por
    # celda se da vuelta (r = 25 en el 51,4 % y r = 1 en el 45,6 %). Ademas 16 + 8 = 24 y no 25:
    # la celda que falta (2 particulas x 10 iteraciones) devuelve r = 14. Lo que caduca es la
    # frecuencia, no el argmax, que sigue en r = 1 dentro del rango publicado.
    '16 de las 25',
    # El libro escribia «(PSO: r=25, m=0,30)» en cinco lugares —resumen, abstract, entrada del
    # indice de figuras, epigrafe de la Figura 10 y conclusion—, como si la optimizacion hubiera
    # hallado los dos hiperparametros. Es lo que H5 niega y lo que el informe dedica once paginas
    # a desmentir: dentro del rango heredado el argmax de la aptitud es r = 1, y el radio lo elige
    # la bateria de evaluacion. Era la ultima inconsistencia ENTRE documentos.
    'PSO: r=25', 'PSO: r = 25',
]
# «NO » en mayusculas hace falta: el proyecto usa el NO enfatico en mayusculas —el cuerpo de la
# lamina 19 dice «La optimizacion NO determina la configuracion adoptada», y las notas del orador
# instruyen con «NO decir «...»» citando textual la frase retirada—. Sin esta variante, cualquier
# frase de RETIRADAS que aparezca citada asi se reportaria como afirmada, o sea un falso positivo.
NEGADORES = ('no ', 'No ', 'NO ', 'ninguna', 'Ninguna', 'ningún', 'sin ', 'tampoco')


def afirmada(texto, frase):
    """True si la frase aparece al menos una vez SIN negacion que la anteceda."""
    for m in re.finditer(re.escape(frase), texto):
        # se mira el fragmento de oracion previo, hasta 90 caracteres o el punto
        ini = max(0, m.start() - 90)
        previo = texto[ini:m.start()]
        previo = previo[previo.rfind('.') + 1:] if '.' in previo else previo
        # el README enfatiza con Markdown («**no** desactiva»), de modo que hay que
        # quitar los marcadores antes de buscar la negacion
        previo = re.sub(r'[*_`]+', ' ', previo)
        if not any(n in previo for n in NEGADORES):
            return True
    return False


for doc in DOCUMENTOS:
    t = DOCUMENTOS[doc]['texto']
    presentes = [s for s in RETIRADAS if afirmada(t, s)]
    negadas = [s for s in RETIRADAS if s in t and not afirmada(t, s)]
    ok(not presentes, f'{doc}: ninguna afirmacion retirada {presentes}')
    if negadas:
        print(f'        (aparecen negadas, correcto: {negadas})')

# Las cifras retiradas hay que buscarlas en LAS DOS notaciones. El abstract en
# ingles del libro sobrevivio trece cifras obsoletas porque esta lista solo miraba
# la coma decimal y el ingles usa punto.
print('\n=== 4b. cifras retiradas, en coma Y en punto ===')
# Solo cifras INEQUIVOCAS. Las de tres decimales del M3FD anterior (0,165 · 0,124 ·
# 0,157 · 0,119 · 0,220 · 0,018 · 0,178 · 0,135) quedan fuera a proposito: colisionan
# con valores vigentes —0,135 es el SD del Top-Hat (0,1352) redondeado— y darian
# falsos positivos. Esas afirmaciones ya las cubre el chequeo de frases del punto 4.
# 0,5781 · 1,7354 · 1,7039 salieron de la lista el 5 de agosto: dejaron de ser
# inequivocas. Las dos primeras coinciden con valores por imagen vigentes de
# pso_por_imagen.csv y la tercera con la aptitud media de las corridas que terminan en
# r = 22 en el estudio de estabilidad. Las afirmaciones que representaban ya las cubre
# el chequeo de frases del punto 4.
CIFRAS_RETIRADAS = ['6,9888', '1,1045', '0,1477', '17,3435', '0,6677', '17,2546',
                    '22,8554', '6,9334', '0,1387', '0,0353', '3,67', '3,44']
for doc in DOCUMENTOS:
    t = DOCUMENTOS[doc]['texto']
    hits = []
    for v in CIFRAS_RETIRADAS:
        for var in (v, v.replace(',', '.')):
            # frontera a la derecha: 3,44 no debe cazar 3,444
            if re.search(re.escape(var) + r'(?!\d)', t):
                hits.append(var)
    ok(not hits, f'{doc}: ninguna cifra retirada, en ninguna notacion {hits[:6]}')

# Las negritas de las tablas se pusieron cuando los valores eran otros y quedaron
# marcando la celda equivocada: en la Tabla 7 el Global señalaba a la piramide de
# Laplace, que es la conclusion que el propio capitulo desmiente.
print('\n=== 4c. negritas de las tablas == optimo de la columna ===')
docx = DOCS / 'Tesis_Borrador_V3.docx'
try:
    import docx as _dx
    rk_t = pd.read_csv(REP / 'ranking_methods.csv', index_col=0)
    ETQ = {'Pirámide de Laplace (LP)': 'PiramideLaplace', 'Ratio Pyramid (RP)': 'RatioPiramide',
           'Wavelet discreta (DWT)': 'DWT', 'DTCWT': 'DTCWT', 'Curvelet (CVT)': 'Curvelet',
           'Top-Hat clásico': 'TopHat_Clasico', 'Propuesta Novedosa': 'Propuesta_Novedosa'}
    ESPERA = {
        ('Método', 'EN', 'SD', 'FE', 'MG', 'MI_vis', 'MI_ir'):
            ('Tabla 4', ['EN', 'SD', 'FE', 'MG', 'MI_vis', 'MI_ir'], dm, 'max'),
        ('Método', 'SF ↑', 'SSIM ↑', 'PSNR ↑'):
            ('Tabla 5', ['SF', 'SSIM', 'PSNR'], dm, 'max'),
        ('Método', 'SD', 'MG', 'SF', 'SSIM', 'PSNR', 'MI_ir', 'Global'):
            ('Tabla 7', ['SD', 'MG', 'SF', 'SSIM', 'PSNR', 'MI_ir', 'avg_rank'], rk_t, 'min'),
    }
    doc_x = _dx.Document(str(docx))
    revisadas = 0
    for tb in doc_x.tables:
        cab = tuple(c.text.strip() for c in tb.rows[0].cells)
        if cab not in ESPERA:
            continue
        nom, cols, src, modo = ESPERA[cab]
        gan = {c: (src[c].idxmax() if modo == 'max' else src[c].idxmin()) for c in cols}
        malas = []
        for fila in tb.rows[1:]:
            cel = list(fila.cells)
            clave = ETQ.get(cel[0].text.strip())
            if clave is None:
                continue
            for j, c in enumerate(cols, start=1):
                neg = any(r.bold for p in cel[j].paragraphs for r in p.runs if r.bold)
                if neg != (gan[c] == clave):
                    malas.append(f'{c}/{clave}')
        ok(not malas, f'{nom}: la negrita marca el optimo {malas[:6]}')
        revisadas += 1
    ok(revisadas == 3, f'se revisaron las 3 tablas con negrita ({revisadas})')
except ImportError:
    print('  AVISO python-docx no instalado: no se reviso la negrita')
    avisos.append('negritas sin revisar (falta python-docx)')

# --------------------------------------------- 5. coherencia entre documentos
print('\n=== 5. coherencia entre documentos ===')
CANTIDADES = {
    'rango de la propuesta': [coma(rk['Propuesta_Novedosa'], 3), punto(rk['Propuesta_Novedosa'], 3),
                              coma(rk['Propuesta_Novedosa'], 2), punto(rk['Propuesta_Novedosa'], 2)],
    'mAP del IR en LLVIP': [coma(ll.loc['IR', 'mAP50'], 3), punto(ll.loc['IR', 'mAP50'], 3)],
    'escenas del conteo M3FD': [str(int(com.loc['VIS', 'escenas']))],
    'tamano del corpus': ['20 pares', '20 pares'],
}
for nombre, variantes in CANTIDADES.items():
    donde = [d for d in DOCUMENTOS if contiene(d, *variantes)]
    ausentes = [d for d in DOCUMENTOS if d not in donde]
    ok(len(donde) >= 2, f'«{nombre}» coincide en {donde} (ausente en {ausentes})',
       blando=True)

# --------------------------------------------- 6. figuras embebidas
print('\n=== 6. figuras embebidas, por md5 ===')
EMBEBIDAS = {
    DOCS / 'Tesis_Borrador_V3.docx': {
        'word/media/image8.png': 'fig_libro_boxplots.png',
        'word/media/image9.png': 'fig_libro_ranking.png',
        'word/media/image13.png': 'fig_libro_propuesta_vs.png',
        'word/media/image14.png': 'fig_libro_pso.png',
        'word/media/image16.png': 'fig_m3fd_detecciones.png',
        # el flujograma no llevaba control y quedo rotulando m = 0,0703 —el optimo de la
        # aptitud paralela— en lugar del peso adoptado, en los DOS entregables
        'word/media/image5.png': 'fig_flujo_propuesta.png',
        'word/media/image2.png': 'fig_morfologia_tophat.png',
        'word/media/image17.png': 'comparacion_aptitudes.png',
        # La curva de aptitud rotulaba «m* = 0,0703» y sombreaba un «rango sugerido m in [0,5-2]»
        # que no es de ningun trabajo; el rango publicado es [0,30; 2,00]. Contradecia a su propio
        # epigrafe, que dice que el optimo dentro del rango cae en su piso. Tenia generador
        # —curva_aptitud_vs_m.py— pero el blob embebido no coincidia con el, asi que no estaba bajo
        # control: es el tercer caso del mismo patron, todos por no estar en esta lista.
        'word/media/image15.png': 'fig_aptitud_vs_m.png',
    },
    DOCS / 'Tesis_Defensa_Presentacion.pptx': {
        'ppt/media/image-8-2.png': 'fig_deck_pso_barrido.png',
        'ppt/media/image-13-1.png': 'fig_deck_llvip_map.png',
        'ppt/media/image-14-1.png': 'fig_deck_m3fd_clases.png',
        'ppt/media/image-15-1.png': 'fig_m3fd_detecciones.png',
        'ppt/media/image-10-1.png': 'cualitativas/montaje_07.png',
        'ppt/media/image-7-1.png': 'fig_flujo_propuesta.png',
        'ppt/media/image-6-1.png': 'fig_morfologia_tophat.png',
        # La figura de la lamina 3 rotulaba «Propuesta (r=25, m=0,070)» —el peso de una corrida
        # descartada, el optimo de la aptitud paralela— y sobrevivio a la correccion del flujograma
        # justamente porque no estaba en esta lista: al no tener gemela en docs/figures, el barrido
        # por md5 nunca la miro. Ahora la produce make_figuras_deck.py::motivacion().
        'ppt/media/image-3-1.png': 'fig_deck_motivacion.png',
    },
}
for contenedor, mapa in EMBEBIDAS.items():
    if not contenedor.exists():
        ok(False, f'falta {contenedor.name}')
        continue
    with zipfile.ZipFile(contenedor) as z:
        nombres = set(z.namelist())
        for interno, fig in mapa.items():
            ruta = DOCS / 'figures' / fig
            if interno not in nombres or not ruta.exists():
                ok(False, f'{contenedor.name}: falta {interno} o {fig}')
                continue
            h1 = hashlib.md5(z.read(interno)).hexdigest()
            h2 = hashlib.md5(ruta.read_bytes()).hexdigest()
            ok(h1 == h2, f'{contenedor.name}: {interno} == {fig}')

# --------------------------------------------- 7. montajes al dia
print('\n=== 7. montajes cualitativos ===')
sys.path.insert(0, str(RAIZ))
try:
    from src.datasets import list_pairs
    pares = list_pairs()
    n_esc = len(pares)
    cual = DOCS / 'figures' / 'cualitativas'
    hay = sorted(cual.glob('montaje_*.png'))
    ok(len(hay) == n_esc, f'hay un montaje por escena ({len(hay)} de {n_esc})')
    # el ultimo par del corpus debe aparecer en la galeria de entrada del informe
    ultimo = pares[-1][0].stem
    penultimo = pares[-2][0].stem
    for doc in ('avances',):
        if doc in DOCUMENTOS:
            ok(contiene(doc, penultimo), f'{doc}: el par {penultimo} figura en la galeria')
except Exception as e:                                             # noqa: BLE001
    ok(False, f'no se pudo comprobar el corpus: {str(e)[:80]}')

# --------------------------------------------- 8. paginacion
print('\n=== 8. paginacion y contadores ===')
for doc, patron, total in (('avances', r'\b(\d{1,3}) */ *(\d{1,3})\b', None),
                           ('deck', r'\b(\d{1,2}) */ *(\d{1,2})\b', 22)):
    if doc not in DOCUMENTOS:
        continue
    with fitz.open(str(DOCUMENTOS[doc]['ruta'])) as d:
        desfasados = []
        for i, pg in enumerate(d, start=1):
            m = re.findall(patron, pg.get_text())
            if total is not None:
                m = [x for x in m if x[1] == str(total)]
            else:
                m = [x for x in m if x[1] == str(d.page_count)]
            if m and int(m[-1][0]) != i:
                desfasados.append((i, m[-1]))
        ok(not desfasados, f'{doc}: contadores sin desfase'
                           + (f' — {desfasados[:4]}' if desfasados else ''))

# --------------------------------------------- 9. desborde del deck
# LibreOffice recorta en el borde de la lamina: lo que no entra no se dibuja y no deja
# rastro en el PDF. Asi la lamina 18 perdio entero el parrafo del conteo por escena, que
# es el resultado de OE5, y nadie lo noto. Comparar el texto de cada shape con el de su
# pagina detecta el recorte sin mirar geometria. Se comparan solo letras y digitos: los
# simbolos (-, °, ∈, ·) no sobreviven igual a las dos extracciones.
print('\n=== 9. texto del deck que no llega al PDF ===')
PPTX = DOCS / 'Tesis_Defensa_Presentacion.pptx'
try:
    from pptx import Presentation

    if 'deck' not in DOCUMENTOS or not PPTX.exists():
        raise FileNotFoundError(PPTX.name)

    def solo_alfanum(s):
        return re.sub(r'[^0-9a-z]+', '', sin_tildes(s or '').lower())

    prs = Presentation(str(PPTX))
    with fitz.open(str(DOCUMENTOS['deck']['ruta'])) as d:
        ok(len(prs.slides) == d.page_count,
           f'el PDF trae una pagina por lamina ({d.page_count} y {len(prs.slides)})')
        perdidos = []
        for i, (lam, pg) in enumerate(zip(prs.slides, d), start=1):
            enpdf = solo_alfanum(pg.get_text())
            for sh in lam.shapes:
                if not sh.has_text_frame:
                    continue
                for par in sh.text_frame.paragraphs:
                    clave = solo_alfanum(par.text)
                    if len(clave) >= 20 and clave not in enpdf:
                        perdidos.append(f'lamina {i} «{sh.name}»')
        ok(not perdidos, 'ningun parrafo se pierde en el recorte'
                         + (f' — {perdidos[:5]}' if perdidos else ''))
except ImportError:
    print('  AVISO python-pptx no instalado: no se reviso el desborde del deck')
    avisos.append('desborde del deck sin revisar (falta python-pptx)')
except FileNotFoundError as e:
    ok(False, f'falta {e} para revisar el desborde del deck', blando=True)

# --------------------------------------------- 10. texto tapado por una figura
# Las figuras del deck son PNG con fondo blanco OPACO. Si el orden de shapes las pone
# despues del cuadro de texto, tapan lo que se solape, y en el PDF el texto sigue estando
# —se puede seleccionar— asi que el chequeo de recorte del bloque 9 no lo ve. Asi la lamina 8
# tenia la segunda linea del titulo debajo del flujograma y la 9 perdia el final de un
# renglon bajo el mapa de calor. Se cruza el bbox de cada span con el de cada imagen.
print('\n=== 10. texto del deck tapado por una figura ===')
UMBRAL_TAPADO = 0.005          # discrimina los solapes reales de los meros roces de borde
if 'deck' in DOCUMENTOS:
    tapados = []
    with fitz.open(str(DOCUMENTOS['deck']['ruta'])) as d:
        for i, pg in enumerate(d, start=1):
            imgs = [fitz.Rect(im['bbox']) for im in pg.get_image_info()]
            if not imgs:
                continue
            for bl in pg.get_text('dict')['blocks']:
                if bl['type'] != 0:
                    continue
                for ln in bl['lines']:
                    for sp in ln['spans']:
                        sr = fitz.Rect(sp['bbox'])
                        if not sp['text'].strip() or sr.get_area() <= 0:
                            continue
                        for ir in imgs:
                            inter = sr & ir
                            if inter.is_empty:
                                continue
                            if inter.get_area() / sr.get_area() >= UMBRAL_TAPADO:
                                tapados.append(f'lamina {i}: {sp["text"].strip()[:40]!r}')
                                break
    ok(not tapados, 'ningun texto queda debajo de una figura'
                    + (f' — {tapados[:5]}' if tapados else ''))

# --------------------------------------------- 8b. la secuencia de pies del informe
# El bloque 8 busca contadores con la forma «N / M», que es la del deck. El informe de
# avances numera con el numero solo, de modo que ese chequeo pasaba sin mirar nada: al
# insertar una pagina en el medio quedaron un pie repetido y otro faltante sin que nadie
# lo notara. Aca se comprueba que la secuencia sea estrictamente consecutiva.
print('\n=== 8b. secuencia de pies del informe de avances ===')
if 'avances' in DOCUMENTOS:
    with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as d:
        pies = []
        for i, pg in enumerate(d, start=1):
            lineas = [l.strip() for l in pg.get_text().splitlines() if l.strip()]
            ult = lineas[-1] if lineas else ''
            pies.append((i, int(ult)) if re.fullmatch(r'\d{1,3}', ult) else (i, None))
    conp = [(i, v) for i, v in pies if v is not None]
    saltos = [f'pag {a[0]}: {a[1]} -> {b[1]}'
              for a, b in zip(conp, conp[1:]) if b[1] != a[1] + 1]
    ok(not saltos, f'los {len(conp)} pies numerados son consecutivos'
                   + (f' — {saltos[:5]}' if saltos else ''))
    rep = [v for _, v in conp]
    ok(len(rep) == len(set(rep)), 'ningun pie se repite'
       + ('' if len(rep) == len(set(rep))
          else f' — repetidos {sorted(v for v in set(rep) if rep.count(v) > 1)}'))
    # Y el encuadre: el pie esta posicionado en absoluto al fondo del div de cada pagina, de
    # modo que si el contenido se pasa de alto el div termina en la pagina SIGUIENTE y se
    # lleva el pie con el. Una pagina sin pie al final es, por tanto, una pagina cuyo
    # contenido se derramo: parte del texto queda huerfano y la siguiente arranca a media
    # frase. La portada es la unica que legitimamente no lleva pie.
    derramadas = [i for i, v in pies if v is None and i > 1]
    ok(not derramadas, 'ninguna pagina derrama su contenido a la siguiente'
                       + (f' — derraman {derramadas}' if derramadas else ''))

# --------------------------------------------- 11. el Excel
# El Excel es un entregable rastreado y no tenia ningun chequeo. Llego a publicar «optimo
# global r* = 25» —lo contrario de H5—, una banda de mAP de una corrida anterior y la
# numeracion vieja de hipotesis. Se revisa en dos planos: las cadenas de texto contra la
# misma lista de afirmaciones retiradas que los otros documentos, y las celdas NUMERICAS
# contra las cifras retiradas, comparando por valor y no por cadena, porque el Excel guarda
# numeros y no texto formateado.
print('\n=== 11. el Excel de tablas ===')
XLSX = DOCS / 'Avances_Tesis_Tablas.xlsx'
try:
    from openpyxl import load_workbook

    if not XLSX.exists():
        raise FileNotFoundError(XLSX.name)
    wb = load_workbook(XLSX, data_only=True)
    textos, numeros = [], []
    for hoja in wb.worksheets:
        for fila in hoja.iter_rows():
            for c in fila:
                if isinstance(c.value, str):
                    textos.append(c.value)
                elif isinstance(c.value, (int, float)) and c.value is not None:
                    numeros.append(float(c.value))
    blob = plano(' '.join(textos))
    ok(len(wb.worksheets) >= 10, f'tiene {len(wb.worksheets)} hojas: {", ".join(wb.sheetnames)}')

    presentes = [s for s in RETIRADAS if afirmada(blob, s)]
    ok(not presentes, f'excel: ninguna afirmacion retirada {presentes}')

    # las cifras retiradas, por valor: se comparan con la tolerancia del formato publicado
    hits = []
    for v in CIFRAS_RETIRADAS:
        try:
            x = float(v.replace(',', '.'))
        except ValueError:
            continue
        nd = len(v.split(',')[1]) if ',' in v else 0
        if any(abs(n - x) < 0.5 * 10 ** (-nd) for n in numeros):
            hits.append(v)
    ok(not hits, f'excel: ninguna cifra retirada en las celdas numericas {hits[:6]}')

    # y las que el Excel publicaba a mano y ya no deben estar
    VIEJAS_EXCEL = ['óptimo global', 'H3 no sostenida', 'r=12', 'm=0.069', '0,957', '0,913']
    quedan = [s for s in VIEJAS_EXCEL if s.lower() in blob.lower()]
    ok(not quedan, f'excel: sin las afirmaciones que publicaba a mano {quedan}')

    # El rango medio no esta como valor: la hoja lo deja en una formula AVERAGE que Excel
    # calcula al abrir, de modo que openpyxl no lo ve. Se verifica la materia prima: que el
    # promedio de los nueve rangos guardados reproduzca el del CSV.
    rk = pd.read_csv(REP / 'ranking_methods.csv', index_col=0)
    hoja = wb['Ranking_Global'] if 'Ranking_Global' in wb.sheetnames else None
    if hoja is None:
        ok(False, 'excel: falta la hoja Ranking_Global')
    else:
        # La fila buscada es la de DATOS, no el subtitulo: este tambien dice «La propuesta»,
        # de modo que hace falta exigir que la fila traiga los nueve rangos. La ultima
        # columna es la formula AVERAGE, que openpyxl lee como None.
        def rangos_de(fila):
            return [c.value for c in fila[1:] if isinstance(c.value, (int, float))]

        fila_prop = next((f for f in hoja.iter_rows()
                          if isinstance(f[0].value, str)
                          and 'PROPUESTA' in f[0].value.upper()
                          and len(rangos_de(f)) == 9), None)
        rangos = rangos_de(fila_prop) if fila_prop else []
        prom = sum(rangos) / len(rangos) if len(rangos) == 9 else None
        esperado = float(rk['avg_rank'].min())
        visto = coma(prom, 3) if prom is not None else 'nada'
        ok(prom is not None and abs(prom - esperado) < 5e-4,
           f'excel: los nueve rangos de la propuesta promedian {visto} '
           f'y el CSV publica {coma(esperado, 3)}')
except ImportError:
    print('  AVISO openpyxl no instalado: no se reviso el Excel')
    avisos.append('excel sin revisar (falta openpyxl)')
except FileNotFoundError as e:
    ok(False, f'falta {e} para revisar el Excel', blando=True)

# El deck sostenia que el PSO devuelve r = 1 «en 16 de las 25 configuraciones» y que las 25
# «convergen al mismo peso». Son cifras del barrido publicado, que corre UNA semilla por celda.
# Con 20 repeticiones por celda —500 corridas— el reparto se da vuelta: r = 25 en el 51,4 % y
# r = 1 en el 45,6 %, y el piso del peso se alcanza en 499 de 500, no en todas.
#
# El bloque 4 no alcanza para vigilarlo, por dos razones medidas:
#  - «que aparece en 8» quedaba SIEMPRE tapada por afirmada(), porque el «no r = 25» que la
#    precede cae en la ventana de 90 caracteres y NEGADORES la anula: falso negativo garantizado.
#  - las notas del orador NO se imprimen en el PDF, de modo que ninguna frase que viva solo ahi
#    puede ser vista por los bloques que leen DOCUMENTOS. La nota de la lamina 19 decia «la
#    busqueda devuelve r = 1», que con las 500 corridas queda literalmente invertida.
# Por eso este bloque va con regex cruda, sin filtro de negacion, y ademas abre el pptx.
print('\n=== 12. formulaciones del barrido PSO retiradas (regex cruda, incluye notas) ===')
FORMULAS_PSO = [
    (r'\b16\s*de\s*(?:las\s*)?25\b', 'r = 1 en 16 de las 25 celdas (barrido de una sola semilla)'),
    (r'que\s+aparece\s+en\s+8\b', 'r = 25 «aparece en 8» (idem, y 16 + 8 = 24, no 25)'),
    # NO se vigila «las 25 configuraciones convergen al mismo peso»: se probo y marcaba tres
    # afirmaciones LEGITIMAS —libro, informe y README— donde la frase habla de las 25 celdas del
    # barrido publicado, y eso es cierto: las 25 tienen m_opt = 0,30, como garantiza el assert de
    # make_figuras_deck.py. En el deck se reescribio por precision, no por error, porque al lado
    # de las 500 corridas «las 25» se volvia ambiguo. Un chequeo asi forzaria a reescribir texto
    # correcto, que es peor que no tenerlo.
    (r'b[uú]squeda\s+devuelve\s+r\s*=\s*1', 'la busqueda «devuelve r = 1» (devuelve r = 25 en el '
                                            '51,4 % y r = 1 en el 45,6 %)'),
    (r'[óo]ptimo\s+de\s+F_o\s+NO\s+es', 'el optimo de F_o «NO es» r = 25 con m = 0,07 (si lo es: '
                                        'es el maximo global de optimo_exacto_fo.csv)'),
]
# el texto del pptx, cuerpo y notas, que es mas amplio que el del PDF
TEXTO_PPTX = ''
if PPTX.exists():
    try:
        from pptx import Presentation as _Pres
        _p = _Pres(str(PPTX))
        _tr = []
        for _sl in _p.slides:
            for _sh in _sl.shapes:
                if _sh.has_text_frame:
                    _tr.append(_sh.text_frame.text)
                if getattr(_sh, 'has_table', False) and _sh.has_table:
                    _tr += [_c.text for _r in _sh.table.rows for _c in _r.cells]
            if _sl.has_notes_slide:
                _tr.append(_sl.notes_slide.notes_text_frame.text)
        TEXTO_PPTX = plano('\n'.join(_tr))
        print(f'  pptx leido: {len(TEXTO_PPTX)} caracteres de cuerpo, tablas y notas')
    except ImportError:
        print('  AVISO python-pptx no instalado: no se revisaron las notas del orador')
        avisos.append('notas del orador sin revisar (falta python-pptx)')

FUENTES_PSO = {k: v['texto'] for k, v in DOCUMENTOS.items()}
if TEXTO_PPTX:
    FUENTES_PSO['pptx (con notas)'] = TEXTO_PPTX
for _pat, _desc in FORMULAS_PSO:
    _rx = re.compile(_pat, re.I)
    _donde = {k: len(_rx.findall(t)) for k, t in FUENTES_PSO.items() if _rx.search(t)}
    ok(not _donde, f'nadie afirma «{_desc}»' + (f' — aparece en {_donde}' if _donde else ''))

# Cruce de citas contra bibliografia: no lo hacia ningun script del repo. verificar_bibliografia.py
# valida el sentido inverso —que lo LISTADO exista en Crossref— y ademas lee un snapshot congelado
# (docs/fuentes/entradas.json), no el listado impreso del libro. Este bloque toma como autoridad el
# capitulo 7 del propio docx, que es la bibliografia que el lector ve.
#
# Motivo real: cuatro citas del proyecto no tenian entrada. Redmon et al. (2016) y Jocher et al.
# (2023), citados en la seccion 12 del informe de avances; y Bai et al. (2015) y Wang et al. (2017)
# en el libro, que quedaron con el anio viejo cuando la auditoria corrigio esas dos entradas a 2012
# y 2014 por DOI.
print('\n=== 13. toda cita en texto tiene entrada en la bibliografia ===')
LIBRO_DOCX = DOCS / 'Tesis_Borrador_V3.docx'
# La segunda letra del apellido puede ser minuscula acentuada (Vázquez, Candès), asi que el token
# es «mayuscula + al menos una letra mas». Las iniciales sueltas («X.», «R. C.») quedan fuera.
TOKEN = r'[A-ZÁÉÍÓÚÑ][\wÁÉÍÓÚÑáéíóúñ\-]+'
ANIO = r'\d{4}|s\.\s?f\.'
NARRATIVA = re.compile(rf'({TOKEN})(?:\s+et\s+al\.|\s+y\s+{TOKEN})?\s*\(({ANIO})')
PARENTETICA = re.compile(rf'\(({TOKEN})(?:\s+et\s+al\.)?,\s*({ANIO})')
# Falsos positivos medidos, con el motivo. No son citas.
IGNORAR = {('cvpr', '2016'), ('cvpr', '2022'), ('icip', '2003'), ('iccv', '2021'),
           ('figura', 's.f.'), ('tabla', 's.f.')}


def _clave(apellido, anio):
    return (sin_tildes(apellido).lower(), re.sub(r'\s', '', anio).lower())


def _autoridad_de(entradas):
    """Todos los apellidos de cada entrada, con su anio: «Ortega y Espinoza (2025)» cita por el
    segundo apellido, y «Bai et al. (2012)» por el primero de tres."""
    aut = set()
    for e in entradas:
        cab, _, resto = e.partition('(')
        m = re.match(rf'\s*({ANIO})', resto)
        if not m:
            continue
        for tok in re.findall(TOKEN, cab):
            aut.add(_clave(tok, m.group(1)))
    return aut


def _citas_de(texto):
    return {_clave(a, y) for rx in (NARRATIVA, PARENTETICA) for a, y in rx.findall(texto)}


try:
    import docx as _docx
    _P = [p.text.strip() for p in _docx.Document(str(LIBRO_DOCX)).paragraphs]
    # La ULTIMA aparicion: la primera es la entrada del indice, y cortar por ella deja el
    # listado vacio sin que nada falle (se comprobo: da 0 entradas y el bloque no ve nada).
    _ir = max(i for i, t in enumerate(_P) if t.startswith('7. REFERENCIAS'))
    _ia = max(i for i, t in enumerate(_P) if t.startswith('8. AP'))
    ENTRADAS_LIBRO = [t for t in _P[_ir + 1:_ia] if t]
    AUT_LIBRO = _autoridad_de(ENTRADAS_LIBRO)
    ok(len(ENTRADAS_LIBRO) >= 36, f'el capitulo 7 del libro lista {len(ENTRADAS_LIBRO)} entradas')

    # El README declara su propia seleccion de referencias, con una entrada —Daubechies (1992)—
    # que no esta entre las del libro. Para el README la autoridad es la union de las dos.
    ENTRADAS_README = []
    _mr = re.search(r'##\s*\d+\.\s*Referencias(.*?)(?=\n##\s|\Z)',
                    (RAIZ / 'README.md').read_text(encoding='utf-8'), re.S)
    if _mr:
        ENTRADAS_README = [re.sub(r'^[-*]\s*', '', l).strip()
                           for l in _mr.group(1).splitlines()
                           if l.strip().startswith(('-', '*'))]
    AUT_README = AUT_LIBRO | _autoridad_de(ENTRADAS_README)
    print(f'  autoridad: {len(AUT_LIBRO)} pares (apellido, anio) del libro '
          f'+ {len(AUT_README - AUT_LIBRO)} propios del README')

    # El cuerpo del libro se lee del docx y se corta ANTES del capitulo 7: si no, las entradas de
    # la bibliografia se cuentan a si mismas como citas y el cruce nunca detecta nada.
    CUERPOS = {'libro': plano('\n'.join(_P[:_ir]))}
    for _d in ('deck', 'avances', 'readme'):
        if _d in DOCUMENTOS:
            CUERPOS[_d] = DOCUMENTOS[_d]['texto']
    for _d, _t in CUERPOS.items():
        _aut = AUT_README if _d == 'readme' else AUT_LIBRO
        _huerf = sorted(c for c in _citas_de(_t) - _aut if c not in IGNORAR)
        ok(not _huerf, f'{_d}: las {len(_citas_de(_t))} citas tienen entrada'
                       + (f' — SIN entrada: {_huerf}' if _huerf else ''))
except ImportError:
    print('  AVISO python-docx no instalado: no se cruzaron las citas')
    avisos.append('citas sin cruzar (falta python-docx)')
except (ValueError, FileNotFoundError) as e:
    ok(False, f'no se pudo delimitar la bibliografia del libro: {e}')

# Todo CSV con una columna de imagen tiene que cubrir EXACTAMENTE el corpus vivo. El defecto que
# motiva el bloque: pso_por_imagen_libre.csv tenia 19 de los 20 pares —le faltaba
# Triclobs_Kaptein_1123, el que sustituyo al corrupto Athena_heather_IR_hei_vis— porque se habia
# corrido antes de la sustitucion. Se publico una tabla que comparaba esa columna contra otra del
# corpus nuevo, y desde el texto no se veia: las dos columnas parecen homologas.
#
# La fecha del archivo NO sirve para detectarlo: pso_por_imagen.csv es anterior al corrimiento que
# se sospechaba y su corpus estaba bien. Lo unico que discrimina es comparar los nombres contra
# list_pairs(), que es la fuente de verdad del corpus.
print('\n=== 14. los CSV por imagen cubren el corpus vivo ===')
# CSV que NO publica ningun entregable: el desajuste se informa como aviso, no como fallo, pero se
# informa, porque un CSV vencido en results/ es una figura equivocada esperando que alguien lo tome.
NO_PUBLICADOS = {
    'fo_ablacion_per_image.csv': 'lo produce eval_fo_optima.py y no lo consume ningun entregable; '
                                 'ademas indexa las imagenes por numero (1..20) y no por nombre, '
                                 'asi que no se puede rastrear a que par corresponde cada fila',
}
try:
    import sys as _sys
    _sys.path.insert(0, str(RAIZ))
    from src.datasets import list_pairs as _lp
    _VIVO = {Path(str(v)).stem for v, _i in _lp()}
    print(f'  corpus vivo: {len(_VIVO)} pares segun list_pairs()')
    _revisados = 0
    for _csv in sorted(REP.glob('*.csv')):
        if '.bak' in _csv.name:
            continue
        try:
            _d = pd.read_csv(_csv, dtype=str, usecols=lambda c: c.lower() in ('imagen', 'image'))
        except (ValueError, pd.errors.EmptyDataError):
            continue
        if _d.empty or not len(_d.columns):
            continue
        _u = set(_d.iloc[:, 0].dropna())
        if not (_u & _VIVO):
            continue      # otro dataset (M3FD, LLVIP) o indexado por numero: no aplica
        _revisados += 1
        _falta, _sobra = sorted(_VIVO - _u), sorted(_u - _VIVO)
        _blando = _csv.name in NO_PUBLICADOS
        ok(not _falta and not _sobra,
           f'{_csv.name}: cubre el corpus vivo'
           + (f' — faltan {_falta} · sobran {_sobra}' if (_falta or _sobra) else ''),
           blando=_blando)
    ok(_revisados >= 6, f'se revisaron {_revisados} CSV indexados por nombre de imagen')
    for _n, _por in NO_PUBLICADOS.items():
        if (REP / _n).exists():
            print(f'        ({_n}: {_por})')
except ImportError as _e:
    print(f'  AVISO no se pudo importar list_pairs ({_e}): corpus sin verificar')
    avisos.append('corpus de los CSV sin verificar (no se pudo importar list_pairs)')

# Los rotulos de las tablas del informe, en orden de aparicion. El defecto que motiva el bloque:
# al agregar una tabla nueva se le puso «Tabla 3f», que ya existia mas adelante, de modo que el
# documento tenia dos Tabla 3f y ademas la 3e quedaba DESPUES de la 3f. Y de antes venia la tabla
# del detector rotulada 10 apareciendo cuatro paginas antes que la 9.
#
# Nada de eso lo veia ningun chequeo: el generador arma cada rotulo a mano en su literal HTML, asi
# que dos tablas pueden llevar el mismo numero sin que nada proteste. Un lector que busca «la
# Tabla 3f» encuentra dos.
print('\n=== 15. rotulos de tabla del informe: sin repetidos y en orden ===')
if 'avances' in DOCUMENTOS:
    # Un EPIGRAFE es «Tabla N.» al principio de su parrafo. Una REFERENCIA EN TEXTO —«la misma
    # advertencia de la Tabla 10.»— tiene la misma forma y no es un rotulo. La primera version de
    # este bloque no las distinguia y dio un falso «Tabla 10 repetida» en cuanto una reescritura
    # cambio los dos puntos de una referencia por un punto. Lo que discrimina es lo que viene ANTES:
    # una referencia siempre va precedida de articulo o preposicion.
    _ANTES_REF = re.compile(r'\b(de|en|la|las|el|los|ver|segun|cf)\s+(la\s+)?$', re.I)
    _rot = []
    with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as _d:
        for _i in range(_d.page_count):
            _t = plano(_d[_i].get_text())
            # «Cuadro» tambien: los anexos rotulan «Cuadro A1.» … «Cuadro A20.» y NINGUN chequeo del
            # repositorio los miraba. Es la peor forma de quedarse corto, porque invita a esquivar el
            # control renombrando: bastaba rotular una tabla «Cuadro» para sacarla del inventario.
            # La zona ('T' o 'C') separa las dos familias, que numeran cada una por su cuenta.
            for _m in re.finditer(r'(Tabla|Cuadro) (A?\d+)([a-z]?)\.', _t):
                if _ANTES_REF.search(_t[max(0, _m.start() - 14):_m.start()]):
                    continue      # referencia en prosa, no rotulo
                _z = _m.group(1)[0] + ('A' if _m.group(2).startswith('A') else '')
                _rot.append((_i + 1, int(_m.group(2).lstrip('A')), _m.group(3),
                             _m.group(0)[:-1], _z))
    print(f'  {len(_rot)} rotulos: ' + ' · '.join(f'p{p}:{e}' for p, _n, _s, e, _z in _rot))

    _etqs = [e for _p, _n, _s, e, _z in _rot]
    _rep = sorted({e for e in _etqs if _etqs.count(e) > 1})
    ok(not _rep, 'ningun rotulo de tabla se repite' + (f' — repetidos {_rep}' if _rep else ''))

    # La monotonia se exige DENTRO DE CADA ZONA y no sobre la lista entera. Las «Tabla N» del cuerpo y
    # las «Cuadro AN» de los anexos son dos series independientes que se intercalan en el documento:
    # compararlas juntas daria un retroceso falso en cada frontera. Lo que no se relaja es la
    # monotonia: cada serie sigue teniendo que ir creciendo.
    _baja = []
    for _zona in sorted({_z for *_r, _z in _rot}):
        _sec = [(p, n, e) for p, n, _s, e, z in _rot if z == _zona]
        _baja += [(_sec[i - 1][2], _sec[i][2]) for i in range(1, len(_sec))
                  if _sec[i][1] < _sec[i - 1][1]]
    ok(not _baja, f'la parte numerica no retrocede dentro de cada serie de rotulos'
                  + (f' — retrocede en {_baja}' if _baja else ''))

    # Dentro de un mismo numero, las letras tienen que ir creciendo. No se exige que la base «N»
    # preceda a «Na»: el informe usa desde antes la convencion inversa —la Tabla 2a esta en la
    # p. 12 y la Tabla 2 en la p. 14— y forzarla obligaria a renumerar todo el documento.
    _mal_letra = []
    for _zona in sorted({_z for *_r, _z in _rot}):
        for _num in sorted({n for _p, n, _s, _e, z in _rot if z == _zona}):
            _ls = [s for _p, n2, s, _e, z in _rot if z == _zona and n2 == _num and s]
            if _ls != sorted(_ls):
                _mal_letra.append((_zona, _num, _ls))
    ok(not _mal_letra, 'dentro de cada número las letras van en orden'
                       + (f' — desordenadas {_mal_letra}' if _mal_letra else ''))
else:
    print('  AVISO no esta el informe de avances: rotulos sin revisar')

# El informe escribe con coma decimal en la prosa, pero varias TABLAS y los rotulos de una figura
# salian con punto: en la misma pagina el mismo numero aparecia de las dos formas —«6,986» en la
# celda y «6.986» en el renglon de abajo—. Eran unas 106 cifras entre la tabla de medias, la del
# barrido, la de Friedman, los rotulos de la figura de cuatro metricas y la lectura de la tabla de
# medias.
#
# El chequeo tiene que excluir TRES cosas que llevan punto legitimamente, y distinguirlas es todo
# el trabajo:
#   1. los separadores de miles —90.000 evaluaciones, 3.012.018 parametros—, que en castellano van
#      con punto;
#   2. los numeros de seccion y de version —§5.8.2, 4.3 Respuestas, Ultralytics 8.4.68—;
#   3. la tabla de configuracion del detector, que DECLARA transcribir los valores literales del
#      archivo args.yaml «con punto decimal, para que puedan copiarse tal cual». Ahi el punto es
#      correcto y cambiarlo romperia el proposito de la tabla.
print('\n=== 16. coma decimal en el informe (salvo lo que legitimamente lleva punto) ===')
if 'avances' in DOCUMENTOS:
    # separador de miles: grupos de exactamente tres digitos, posiblemente encadenados
    # OJO: la forma no alcanza para distinguir el separador de miles del decimal de tres cifras.
    # «6.840» encaja igual de bien en \d{1,3}(\.\d{3})+ que «90.000», y la primera version de este
    # bloque daba por buenas las 65 cifras de la tabla de medias por eso. Lo que sí discrimina en
    # este documento es lo que viene DESPUES: los separadores de miles siempre preceden a un
    # sustantivo en minuscula —«90.000 evaluaciones», «3.012.018 parametros», «2.000 imagenes»—
    # mientras los decimales de una tabla van seguidos de otro numero o de fin de celda.
    _MILES = re.compile(r'^\d{1,3}(\.\d{3})+$')
    # Solo nombres de archivo, versiones y enlaces. NO se excluye por la cercania de «Tabla» o
    # «Figura»: la primera version del chequeo lo hacia y tapaba las primeras celdas de cada
    # tabla, que es justo donde estaban la mitad de los defectos —contaba 43 donde habia mas de
    # cien—. Un chequeo que excluye de mas es peor que no tenerlo.
    _CTX_OK = re.compile(r'\.py|\.csv|\.json|\.pt|\.yaml|doi|http|§|Ultralytics|PyTorch|cu\d|'
                         r'yolov', re.I)
    # Los numeros de subseccion —«4.1 Elementos estructurantes»— no son decimales. En lugar de
    # adivinarlos con una regex se leen del HTML que el generador escribe al lado del PDF, que es
    # donde los encabezados estan marcados como tales. Si el HTML no esta, el conjunto queda
    # vacio y el chequeo se pone mas estricto, no menos.
    _SECS = set()
    _html = DOCS / '_local' / 'Avances_Tesis.html'
    if _html.exists():
        _h = _html.read_text(encoding='utf-8', errors='replace')
        for _m in re.finditer(r'<h[23][^>]*>\s*(\d+(?:\.\d+)+)', _h):
            _SECS.add(_m.group(1))
        print(f'  {len(_SECS)} numeros de subseccion excluidos: {sorted(_SECS)}')
    _con_punto = []
    with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as _d:
        for _i in range(_d.page_count):
            _t = plano(_d[_i].get_text())
            # la pagina que declara el punto literal queda fuera por su propia declaracion
            if 'con punto decimal' in _t:
                continue
            for _m in re.finditer(r'\d+\.\d+(\.\d+)*', _t):
                _s = _m.group(0)
                if _MILES.match(_s) and re.match(r'\s+[a-záéíóúñ]', _t[_m.end():_m.end() + 3]):
                    continue
                # numero de subseccion, y ademas seguido del titulo en mayuscula
                if _s in _SECS and re.match(r'\s+[A-ZÁÉÍÓÚÑ]', _t[_m.end():_m.end() + 3]):
                    continue
                _ctx = _t[max(0, _m.start() - 34):_m.end() + 20]
                if _CTX_OK.search(_ctx):
                    continue
                _con_punto.append((_i + 1, _s, _ctx.strip()))
    ok(not _con_punto,
       f'ninguna cifra con punto decimal fuera de la tabla que lo declara'
       + (f' — {len(_con_punto)}: ' + '; '.join(f'p{p} «{s}»' for p, s, _c in _con_punto[:6])
          if _con_punto else ''))
    if _con_punto:
        for _p, _s, _c in _con_punto[:4]:
            print(f'        p.{_p}: ...{_c}...')
else:
    print('  AVISO no esta el informe de avances: coma decimal sin revisar')

# El informe llevaba 83 paginas sin indice y sin un solo marcador: en la reunion no habia forma de
# ir a una seccion salvo desplazandose. Los marcadores se escriben en un posproceso con PyMuPDF
# —Edge headless no los emite, y la regla CSS bookmark-level solo la entiende WeasyPrint—, de modo
# que son un paso aparte que puede fallar en silencio: si el posproceso no corre, el PDF sale igual
# y nadie se entera.
print('\n=== 17. indice y marcadores del informe ===')
if 'avances' in DOCUMENTOS:
    with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as _d:
        _toc = _d.get_toc()
        _npg = _d.page_count
        # El umbral era «>= 40» con 84 marcadores en el documento: se podia perder la mitad del
        # informe y el chequeo pasaba. Ahora se exige uno por pagina salvo la portada y el indice,
        # que es el invariante real y el que el bloque 26 usa del otro lado.
        ok(len(_toc) == _npg - 2,
           f'el PDF trae un marcador por pagina salvo portada e indice ({len(_toc)} de {_npg - 2})')

        # Cada marcador tiene que caer en la pagina cuyo PIE IMPRESO es ese mismo numero. Si el
        # posproceso se corriera sobre un PDF con otra paginacion, los marcadores quedarian
        # corridos y el indice seguiria pareciendo correcto.
        def _pie_de(_p):
            _ls = [l.strip() for l in _d[_p - 1].get_text().splitlines() if l.strip()]
            return _ls[-1] if _ls else ''

        _fuera = [(t, p) for _n, t, p in _toc if not 1 <= p <= _npg]
        _corridos = [(t, p) for _n, t, p in _toc if 1 <= p <= _npg and _pie_de(p) != str(p)]
        ok(not _fuera, 'ningun marcador apunta fuera del documento'
                       + (f' — {_fuera[:3]}' if _fuera else ''))
        ok(not _corridos, f'los {len(_toc)} marcadores caen en la pagina de su pie impreso'
                          + (f' — corridos {[(t[:30], p) for t, p in _corridos[:3]]}'
                             if _corridos else ''))

        # El indice impreso: cada entrada tiene que apuntar a una pagina que exista y su rango
        # tiene que empezar donde arranca la seccion segun los marcadores.
        #
        # La pagina del indice SE BUSCA, no se da por sentada. Antes esto leia _d[1] a secas, o sea
        # la segunda pagina fisica, y al meterle la carilla de resumen al frente el chequeo empezo a
        # leer el resumen y a informar que el indice citaba «1 de los 37 comienzos de seccion». El
        # fallo era del chequeo, no del documento. Buscarlo tampoco lo debilita: si no hubiera
        # ninguna pagina de indice, no hay nada que encontrar y falla igual.
        _pag_idx = next((_i for _i in range(1, min(6, _npg) + 1)
                         if plano(_d[_i - 1].get_text()).lstrip().startswith('Índice')), None)
        ok(_pag_idx is not None, 'hay una pagina de indice entre las cinco primeras')
        _t2 = plano(_d[(_pag_idx or 2) - 1].get_text())

        # Y el orden del frente, que es la razon de ser de la carilla: el resumen tiene que venir
        # ANTES del indice. Si un dia alguien reordena los bloques, el resumen deja de ser lo
        # primero que se lee y pierde el sentido, sin que nada mas lo delate.
        _pag_res = next((_i for _i in range(1, min(6, _npg) + 1)
                         if plano(_d[_i - 1].get_text()).lstrip().startswith('Resumen')), None)
        ok(_pag_res is not None and _pag_idx is not None and _pag_res < _pag_idx,
           f'el resumen (pag. {_pag_res}) va antes del indice (pag. {_pag_idx})')
        _rangos = re.findall(r'(\d{1,3})(?:–(\d{1,3}))?\s*(?=\d{1,3}\.|Anexos|$)', _t2)
        _inicios = sorted({p for _n, _t, p in _toc if _n == 1})
        _citadas = sorted({int(a) for a, _b in re.findall(r'\b(\d{1,3})(?:–(\d{1,3}))?\b', _t2)
                           if int(a) in _inicios})
        ok(len(_citadas) >= 12,
           f'el indice cita {len(_citadas)} de los {len(_inicios)} comienzos de seccion')
else:
    print('  AVISO no esta el informe de avances: indice y marcadores sin revisar')

# Toda imagen embebida sin gemela en docs/figures es una imagen que NO SE PUEDE REHACER: no se sabe
# que script la produjo ni con que datos, de modo que si las cifras cambian nadie se entera. Asi
# sobrevivieron TRES figuras rotulando m = 0,0703 —el optimo de la aptitud paralela, que la tesis no
# adopta— mucho despues de que el flujograma se corrigiera: el flujograma estaba en EMBEBIDAS y
# ellas no, de modo que el barrido por md5 nunca las miro.
#
# Las que quedan son legitimas y estan declaradas una por una con su motivo. El chequeo no exige que
# la lista se vacie: exige que no aparezca una NUEVA sin declarar, que es lo que hay que impedir.
print('\n=== 18. imagenes embebidas sin gemela: todas declaradas ===')
SIN_GEMELA = {
    'Tesis_Borrador_V3.docx': {
        'image1.png': 'logo de la Universidad Comunera; no sale de ningun dato',
        'image3.png': 'flujograma conceptual del Top-Hat clasico (disco unico, m = 1); correcto y '
                      'sin cifras del estudio',
        'image4.png': 'esquema del banco de cinco elementos estructurantes; conceptual',
        'image6.png': 'esquema ilustrativo del PSO. Rotula «gbest (r = 25; m = 0,0703)», que es el '
                      'optimo de la aptitud paralela y NO el peso adoptado; el epigrafe de la '
                      'Figura 5 lo declara. Repintarla exigiria escribirle un generador a una '
                      'figura que no reporta datos',
        'image7.png': 'PENDIENTE: montaje de 3 escenas x 6 entradas con imagenes FUSIONADAS, sin '
                      'generador. Es el unico caso que si deberia tenerlo, porque muestra '
                      'resultados del operador y hoy no se puede rehacer',
    },
    'Tesis_Defensa_Presentacion.pptx': {
        'image-6-2.png': 'ecuacion renderizada', 'image-7-2.png': 'ecuacion renderizada',
        'image-7-3.png': 'ecuacion renderizada', 'image-8-1.png': 'la formula de F_o renderizada',
    },
}
_FIG_MD5 = set()
for _p in (RAIZ / 'docs' / 'figures').rglob('*'):
    if _p.is_file() and _p.suffix.lower() in ('.png', '.jpg', '.jpeg'):
        _FIG_MD5.add(hashlib.md5(_p.read_bytes()).hexdigest())
for _cont, _decl in SIN_GEMELA.items():
    _ruta = DOCS / _cont
    if not _ruta.exists():
        ok(False, f'falta {_cont}')
        continue
    with zipfile.ZipFile(_ruta) as _z:
        _huerf = {n.split('/')[-1] for n in _z.namelist()
                  if '/media/' in n and n.lower().endswith(('.png', '.jpg', '.jpeg'))
                  and hashlib.md5(_z.read(n)).hexdigest() not in _FIG_MD5}
    _nuevas = sorted(_huerf - set(_decl))
    _yano = sorted(set(_decl) - _huerf)
    ok(not _nuevas, f'{_cont}: las {len(_huerf)} imagenes sin gemela estan declaradas'
                    + (f' — SIN DECLARAR: {_nuevas}' if _nuevas else ''))
    if _yano:
        print(f'        ({_yano} ya tienen gemela: se pueden sacar de la lista)')
_pend = [f'{c}/{n}' for c, d in SIN_GEMELA.items() for n, m in d.items()
         if m.startswith('PENDIENTE')]
if _pend:
    print(f'  AVISO {len(_pend)} sin generador y con datos: {_pend}')
    avisos.append(f'figuras con datos y sin generador: {_pend}')

# Todo paquete que el codigo importe tiene que estar en requirements.txt. El defecto que motiva el
# bloque: faltaban CINCO —PyMuPDF, python-docx, python-pptx, Pillow y PyYAML— de modo que un clon
# recien bajado no podia generar el informe ni correr estos verificadores. Fallaba con
# ModuleNotFoundError, y no habia forma de enterarse sin probarlo en una maquina limpia.
print('\n=== 19. requirements.txt declara todo lo que el codigo importa ===')
_STDLIB = set('''os sys re json math time datetime pathlib collections itertools functools argparse
subprocess shutil zipfile hashlib io copy random warnings unicodedata typing dataclasses glob
statistics textwrap html csv traceback base64 xml multiprocessing urllib abc contextlib enum
tempfile string operator pickle gzip struct threading queue logging inspect importlib ast'''.split())
# nombre del modulo -> nombre del paquete en PyPI, cuando no coinciden
_PAQ = {'cv2': 'opencv-python', 'skimage': 'scikit-image', 'PIL': 'Pillow', 'fitz': 'PyMuPDF',
        'docx': 'python-docx', 'pptx': 'python-pptx', 'pywt': 'PyWavelets',
        'sklearn': 'scikit-learn', 'yaml': 'PyYAML'}
_req_txt = (RAIZ / 'requirements.txt').read_text(encoding='utf-8').lower()
_imp = {}
for _py in sorted(list((RAIZ / 'src').rglob('*.py')) + list((RAIZ / 'experiments').rglob('*.py'))):
    for _m in re.finditer(r'^\s*(?:import|from)\s+([A-Za-z_]\w*)',
                          _py.read_text(encoding='utf-8', errors='replace'), re.M):
        _mod = _m.group(1)
        if _mod in _STDLIB or _mod in ('src', 'experiments', '__future__'):
            continue
        # Los MODULOS LOCALES no van en requirements.txt: un script de experiments/ que importa a otro
        # del mismo directorio —«import run_complementariedad_escenas as base», para no tener dos
        # implementaciones del mismo criterio de emparejado— no depende de ningun paquete de PyPI. Sin
        # esta excepcion el chequeo reclamaba que se declararan como dependencias, que es un falso
        # positivo: el modulo esta en el repositorio, al lado del que lo importa.
        if (RAIZ / 'experiments' / f'{_mod}.py').exists() or (RAIZ / f'{_mod}.py').exists():
            continue
        _imp.setdefault(_PAQ.get(_mod, _mod), set()).add(_py.name)
_nodecl = {p: sorted(v)[:3] for p, v in _imp.items() if p.lower() not in _req_txt}
ok(not _nodecl, f'los {len(_imp)} paquetes importados estan en requirements.txt'
                + (f' — SIN DECLARAR: {_nodecl}' if _nodecl else ''))
# y que se puedan importar de verdad en este entorno
_faltan_inst = []
for _mod in ('fitz', 'docx', 'pptx', 'PIL', 'yaml', 'cv2', 'pandas', 'numpy'):
    try:
        __import__(_mod)
    except ImportError:
        _faltan_inst.append(_mod)
ok(not _faltan_inst, 'los modulos criticos se importan en este entorno'
                     + (f' — faltan {_faltan_inst}' if _faltan_inst else ''))

# --------------------------------------------- 23. texto del deck dibujado FUERA de la lamina
# POR QUE, Y EN QUE SE DIFERENCIA DEL BLOQUE 9. El 9 comprueba que el texto de cada shape LLEGUE al
# PDF, y con eso atrapa lo que LibreOffice recorta del todo. Pero hay un caso que deja pasar: el
# texto que SI se dibuja y queda pasado el borde de la lamina. Entra en la capa de texto del PDF —de
# modo que el bloque 9 lo encuentra y da por bueno— y en la proyeccion no esta.
#
# Paso de verdad y estuvo asi mucho tiempo: el ultimo renglon de la prosa de cierre de la lamina 12
# se dibujaba en y = 408,9 pt sobre una lamina de 405, o sea que el tribunal no veia el cierre de la
# lamina de resultados. Se descubrio mirando el render, no con un chequeo.
#
# El margen de 1 pt es tolerancia de redondeo de la extraccion, no holgura de diseño.
print('\n=== 23. texto del deck dentro de los limites de la lamina ===')
if 'deck' in DOCUMENTOS:
    with fitz.open(str(DOCUMENTOS['deck']['ruta'])) as _d:
        _afuera = []
        for _i, _pg in enumerate(_d, start=1):
            _bl = [_b for _b in _pg.get_text('blocks') if _b[4].strip()]
            if not _bl:
                continue
            _abajo = max(_b[3] for _b in _bl)
            _dcha = max(_b[2] for _b in _bl)
            if _abajo > _pg.rect.height - 1 or _dcha > _pg.rect.width - 1:
                _peor = max(_bl, key=lambda _b: _b[3])
                _afuera.append((_i, round(_abajo, 1), round(_pg.rect.height, 1),
                                re.sub(r'\s+', ' ', _peor[4]).strip()[-40:]))
        ok(not _afuera, f'en las {_d.page_count} laminas ningun texto se sale del borde'
                        + (f' — {_afuera[:3]}' if _afuera else ''))

# --------------------------------------------- 22. la carilla de resumen no introduce cifras nuevas
# POR QUE. La pagina 2 es un resumen en lenguaje llano, y es lo primero que lee el director. Es
# tambien el lugar donde una cifra mal copiada hace mas daño: quien la lee ahi no va a ir a buscarla
# a la pagina 36. La regla es que el resumen NO PUEDE ESTRENAR NINGUNA CIFRA. Todo decimal que cite
# tiene que aparecer, con ese mismo redondeo, en alguna otra pagina del informe, que es donde esta
# la tabla y el metodo que lo produce. Si el resumen dijera 0,097 donde el cuerpo dice 0,906, esto
# falla; si el cuerpo cambia y el resumen queda viejo, tambien.
print('\n=== 22. las cifras de la carilla de resumen estan respaldadas en el cuerpo ===')
if 'avances' in DOCUMENTOS:
    with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as _d:
        _pag_res = next((_i for _i in range(1, min(6, _d.page_count) + 1)
                         if plano(_d[_i - 1].get_text()).lstrip().startswith('Resumen')), None)
        if _pag_res is None:
            ok(False, 'no se encontro la carilla de resumen')
        else:
            _t_res = _d[_pag_res - 1].get_text()
            _cuerpo_txt = "\n".join(_d[_i].get_text() for _i in range(_d.page_count)
                                    if _i != _pag_res - 1)
            # los decimales con coma son las afirmaciones; los enteros son conteos y remisiones
            _dec_res = sorted(set(re.findall(r'\d+,\d+', _t_res)))
            _huerfanas = [x for x in _dec_res if x not in _cuerpo_txt]
            ok(not _huerfanas,
               f'las {len(_dec_res)} cifras decimales del resumen {_dec_res} aparecen en el cuerpo'
               + (f' — SIN RESPALDO: {_huerfanas}' if _huerfanas else ''))
            # y las remisiones de pagina tienen que existir y llevar a la seccion que nombran
            _rem = [(int(_s), int(_p)) for _s, _p in
                    re.findall(r'la (?:secci[oó]n )?(\d{1,2}) \(p[aá]g\. (\d{1,3})\)', _t_res)]
            _rem += [(int(_s), int(_p)) for _s, _p in
                     re.findall(r'secci[oó]n (\d{1,2}) \(p[aá]g\. (\d{1,3})\)', _t_res)]
            _malas = []
            for _sec, _p in set(_rem):
                if not 1 <= _p <= _d.page_count:
                    _malas.append((_sec, _p, 'fuera del documento'))
                    continue
                _prim = next((_l.strip() for _l in _d[_p - 1].get_text().splitlines()
                              if _l.strip()), '')
                if not _prim.startswith(f'{_sec}.'):
                    _malas.append((_sec, _p, f'ahi arranca «{_prim[:34]}»'))
            ok(_rem and not _malas,
               f'las {len(set(_rem))} remisiones del resumen caen en la seccion que nombran'
               + (f' — {_malas[:4]}' if _malas else ''))

# --------------------------------------------- 26. inventario del informe: ninguna pagina se pierde
# POR QUE, Y ES EL CHEQUEO MAS IMPORTANTE DE LOS QUE MIRAN EL INFORME. Los pies del informe se
# escribian a mano, uno por pagina, y el generador abortaba si la secuencia no era 2..N. Eso tenia un
# efecto lateral valiosisimo que nadie habia declarado: si un bloque de pagina se perdia en una
# edicion, quedaba un hueco en la secuencia y el generador no escribia nada. Al pasar pie() a un
# contador posicional —para poder mover paginas sin renumerar a mano— ese efecto lateral desaparecio:
# perder un bloque ahora produce un documento de 85 paginas consecutivas, con el indice y los
# marcadores coherentes, y TODO pasa en verde. Se cambio una clase de error por otra.
# Este bloque es la reposicion. El HTML que el generador deja en docs/_local es el molde del PDF, asi
# que las dos cosas tienen que cuadrar exactamente; y los titulos de la seccion 5 se declaran uno por
# uno, porque son el contenido que se movio a un anexo y «nada se elimino» tiene que ser verificable y
# no una afirmacion. Los marcadores son los bloques menos DOS: la portada, que no lleva ninguno, y la
# pagina del indice, que se inserta despues de armar el arbol y por eso queda sin marcador propio.
print('\n=== 26. inventario del informe: el PDF tiene todas las paginas del HTML que lo produjo ===')
_html_inf = RAIZ / 'docs' / '_local' / 'Avances_Tesis.html'
if 'avances' in DOCUMENTOS and _html_inf.exists():
    _h = _html_inf.read_text(encoding='utf-8')
    # Con regex y no con split de cadena: la portada es <div class="page portada">, y partir por
    # '<div class="page"' la saltea. La primera version de este bloque lo hacia asi y contaba 85
    # bloques contra 86 paginas del PDF, informando un descuadre que no existia.
    _divs = re.split(r'<div class="page[^"]*"', _h)[1:]
    _tits = []
    for _b in _divs:
        _mh = re.search(r'<h2>(.*?)</h2>', _b, re.S)
        _mp = re.search(r'<div class="pie">(\d+)</div>', _b)
        if _mh and _mp:
            _tits.append((re.sub(r'\s+', ' ', re.sub(r'<[^>]+>', '', _mh.group(1))).strip(),
                          int(_mp.group(1))))
    with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as _d:
        _npg_pdf, _ntoc = _d.page_count, len(_d.get_toc())
        _primera = {}
        for _i in range(_npg_pdf):
            _ls = [l.strip() for l in _d[_i].get_text().splitlines() if l.strip()]
            _primera[_i + 1] = _ls[0] if _ls else ''
    ok(_npg_pdf == len(_divs),
       f'el PDF tiene {_npg_pdf} paginas y el HTML {len(_divs)} bloques: coinciden'
       if _npg_pdf == len(_divs) else
       f'el PDF tiene {_npg_pdf} paginas y el HTML {len(_divs)} bloques: NO coinciden')
    ok(_ntoc == len(_divs) - 2,
       f'los marcadores son {_ntoc} = {len(_divs)} bloques menos la portada y el indice'
       + ('' if _ntoc == len(_divs) - 2 else f' — esperaba {len(_divs) - 2}'))
    # cada titulo del HTML tiene que abrir la pagina que su pie declara. Un bloque que Edge fusione
    # o parta se ve aca y en ningun otro lado.
    _desfasados = [(t, p, _primera.get(p, '')[:40]) for t, p in _tits
                   if not _primera.get(p, '').startswith(t[:34])]
    ok(not _desfasados,
       f'los {len(_tits)} titulos del HTML abren la pagina que declara su pie'
       + (f' — {_desfasados[:3]}' if _desfasados else ''))

    # LOS DOCE TITULOS DE LA SECCION 5, declarados. El valor dice donde tiene que estar cada uno:
    # 'seccion' en el cuerpo y 'anexo' despues de la ultima seccion numerada. Mover un bloque cambia
    # su valor aca y nada mas; PERDER un bloque falla, que es el punto.
    _S5 = {
        'barrido de configuraciones': 'seccion',
        'convergencia y óptimo': 'seccion',
        # las diez que pasaron al anexo A21. Cambiar un valor de aqui es TODO lo que hace falta si
        # algun dia una vuelve al cuerpo; perder un bloque, en cambio, falla y lo nombra.
        'justificación del peso adoptado': 'anexo',
        'la equivalencia del realce físico': 'anexo',
        'rango dinámico y tensión de criterios': 'anexo',
        'estabilidad del barrido en': 'anexo',
        'las 500 corridas, una por una': 'anexo',
        'el óptimo exacto, por enumeración': 'anexo',
        'por qué el barrido de la referencia dispersa': 'anexo',
        'el mismo barrido, imagen por imagen': 'anexo',
        'el mismo barrido con el peso libre': 'anexo',
        'el registro de las 500 corridas': 'anexo',
        # y las dos paginas nuevas que las reemplazan en la seccion 5
        'el punto de operación adoptado': 'seccion',
        'la conclusión de H5': 'seccion',
    }
    # la frontera entre el cuerpo y los anexos se DERIVA: la ultima pagina con un titulo que empieza
    # por un digito. Escribirla a mano la dejaria vieja al primer cambio de tamaño del informe.
    _ult_cuerpo = max([p for t, p in _tits if re.match(r'\d+\.', t)] or [0])
    _falt, _mal_lugar, _dup = [], [], []
    for _frag, _donde in _S5.items():
        _hits = [p for t, p in _tits if _frag in t]
        if not _hits:
            _falt.append(_frag)
        elif len(_hits) > 1:
            _dup.append((_frag, _hits))
        else:
            _p = _hits[0]
            if _donde == 'anexo' and _p <= _ult_cuerpo:
                _mal_lugar.append((_frag, _p, 'deberia estar en el anexo'))
            if _donde == 'seccion' and _p > _ult_cuerpo:
                _mal_lugar.append((_frag, _p, 'deberia estar en el cuerpo'))
    ok(not _falt, f'los {len(_S5)} bloques declarados de la seccion 5 siguen en el documento'
                  + (f' — PERDIDOS: {_falt}' if _falt else ''))
    ok(not _dup, 'ninguno de esos bloques aparece dos veces'
                 + (f' — {_dup}' if _dup else ''))
    ok(not _mal_lugar, f'cada uno esta del lado que declara (cuerpo hasta la pag. {_ult_cuerpo})'
                       + (f' — {_mal_lugar[:4]}' if _mal_lugar else ''))
else:
    print('  AVISO no esta el HTML del informe en docs/_local: inventario sin revisar')

# --------------------------------------------- 25. el README no declara un conteo de paginas viejo
# POR QUE. El README es lo que lee quien clona el repositorio, y describe el estado ACTUAL: cuando
# dice «Avances_Tesis.pdf — Informe de avances (85 págs)» y el informe tiene 86, la primera cifra que
# ve un tercero ya esta mal. Y no es una cifra que alguien vaya a revisar: se escribio una vez y se
# quedo. Aparecio al reestructurar la seccion 5, que cambia el conteo, pero estaba desactualizada
# desde antes. Se controla solo el README a proposito: ESTADO_Y_PENDIENTES.md es un registro con
# entradas fechadas, y ahi «85 paginas» bajo un encabezado de agosto es verdadero para su fecha.
print('\n=== 25. el README declara el conteo de paginas real de cada entregable ===')
if 'readme' in DOCUMENTOS:
    _txt_rm = (RAIZ / 'README.md').read_text(encoding='utf-8')
    _NOMBRE = {'avances': 'Avances_Tesis.pdf', 'libro': 'Tesis_Borrador_V3.pdf',
               'deck': 'Tesis_Defensa_Presentacion.pdf'}
    _mal = []
    for _cl, _arch in _NOMBRE.items():
        if _cl not in DOCUMENTOS:
            continue
        _real = DOCUMENTOS[_cl]['paginas']
        for _m in re.finditer(re.escape(_arch), _txt_rm):
            # la ventana es el resto del renglon: ahi va la descripcion del archivo
            _fin = _txt_rm.find('\n', _m.end())
            _vent = _txt_rm[_m.end():_fin if _fin > 0 else len(_txt_rm)]
            for _d in re.findall(r'(\d{2,3})\s*(?:p[áa]gs?\b|p[áa]ginas)', _vent):
                if int(_d) != _real:
                    _mal.append(f'{_arch} dice {_d} y tiene {_real}')
    ok(not _mal, f'los conteos de paginas del README coinciden con los PDF'
                 + (f' — {_mal}' if _mal else ''))

# --------------------------------------------- 24. el resumen cita el mAP VIGENTE, no uno que existio
# POR QUE. El bloque 22 exige que toda cifra del resumen APAREZCA en el cuerpo, y con eso alcanzaba
# hasta que LLVIP paso a cinco semillas. Ahi el chequeo se quedo corto de la forma mas incomoda: el
# resumen seguia diciendo que la propuesta llega a «0,906» y que hay «un solo entrenamiento por
# entrada», y 22 lo dejaba pasar porque 0,906 SI aparece en el cuerpo — la seccion 13 lo cita para
# explicar que era el valor de la semilla publicada y el mas bajo de los cinco. O sea que el bloque
# verificaba PROCEDENCIA y no VIGENCIA: una cifra retirada, que el cuerpo menciona justamente para
# desautorizarla, seguia sirviendo de respaldo. El resultado era un informe que en la pagina 2 decia
# una cosa y en la 53 la contraria, y la pagina 2 es la que el director lee primero.
# COMO SE ARREGLA. Todo mAP que el resumen le atribuya a una entrada del detector tiene que ser una
# MEDIA del estudio de semillas, con su mismo redondeo. Y el resumen no puede declarar una sola
# semilla en LLVIP cuando el estudio tiene cinco.
print('\n=== 24. los mAP del resumen son las medias vigentes del estudio de semillas ===')
_sem_res = REP / 'semillas_llvip_resumen.csv'
if 'avances' in DOCUMENTOS and _sem_res.exists():
    _sm = pd.read_csv(_sem_res, index_col=0)
    _n_sem = int(pd.read_csv(REP / 'detection_llvip_semillas.csv').semilla.nunique())
    with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as _d:
        _i_res = next((_i for _i in range(min(6, _d.page_count))
                       if plano(_d[_i].get_text()).lstrip().startswith('Resumen')), None)
    if _i_res is None:
        ok(False, 'no se encontro la carilla de resumen')
    else:
        with fitz.open(str(DOCUMENTOS['avances']['ruta'])) as _d:
            _t = re.sub(r'\s+', ' ', _d[_i_res].get_text())
        # el parrafo del hallazgo de deteccion: de «no se traslada a la tarea» hasta el punto final
        _m = re.search(r'no se traslada a la tarea\.(.{0,900}?)(?=Y el conjunto de nueve)', _t)
        if not _m:
            ok(False, 'no se encontro en el resumen el parrafo del hallazgo de deteccion')
        else:
            _frag = _m.group(1)
            # los mAP validos: las medias de las nueve entradas, a tres y a cuatro decimales
            _validos = set()
            for _col in ('mAP50_media', 'mAP50_95_media'):
                for _v in _sm[_col]:
                    _validos |= {f'{_v:.3f}'.replace('.', ','), f'{_v:.4f}'.replace('.', ',')}
            _desv = {f'{_v:.4f}'.replace('.', ',') for _v in _sm.mAP50_desv}
            _desv |= {f'{_sm.mAP50_desv.median():.4f}'.replace('.', ',')}
            _citados = set(re.findall(r'0,\d{3,4}', _frag))
            _viejos = sorted(_citados - _validos - _desv)
            ok(not _viejos,
               f'los {len(_citados)} mAP que el resumen atribuye a LLVIP son medias del estudio'
               + (f' — RETIRADOS: {_viejos}' if _viejos else ''))
            _una = [_f for _f in ('un solo entrenamiento', 'una sola semilla', 'un único entrenamiento')
                    if _f in plano(_frag)]
            ok(not _una or _n_sem == 1,
               f'el resumen no declara un unico entrenamiento habiendo {_n_sem} semillas'
               + (f' — dice «{_una[0]}»' if _una else ''))
            ok(f'{_n_sem} veces' in _frag or f'{_n_sem} semillas' in _frag
               or f'{_n_sem} entrenamientos' in _frag,
               f'el resumen declara las {_n_sem} semillas al dar el hallazgo de deteccion')

        # Y LAS DOS UNIDADES CON QUE DA SUS RESULTADOS SE EXPLICAN EN LA PROPIA CARILLA. El resumen
        # esta escrito en lenguaje llano, pero llegaba con «3,39 de rango medio» y «0,961 de mAP» sin
        # decir en ninguna parte que son. El rango medio es el peor de los dos: es un promedio de
        # PUESTOS, asi que mas bajo es mejor, y sin eso «primero de siete con 3,39» se lee al reves.
        # La lista de glosas admitidas puede crecer si algun dia se reformula: lo que se exige es que
        # haya alguna, no una redaccion en particular.
        _GLOSAS = {'mAP': ('puntaje del detector', 'entre 0 y 1', 'de 0 a 1'),
                   'rango medio': ('promedio del puesto', 'promedio de los puestos',
                                   'más bajo es mejor', 'mas bajo es mejor')}
        _sin = [_u for _u, _gs in _GLOSAS.items()
                if _u.lower() in _t.lower() and not any(_g in plano(_t) for _g in map(plano, _gs))]
        ok(not _sin,
           f'las {len(_GLOSAS)} unidades con que el resumen da sus resultados se explican ahi mismo'
           + (f' — SIN GLOSA: {_sin}' if _sin else ''))

# --------------------------------------------- 21. los nombres de escena cubren el corpus vivo
# POR QUE. Los generadores traducen el nombre de archivo de cada par a un rotulo legible con un
# diccionario a mano: «Triclobs_Bosnia_R» -> «Bosnia». Si al diccionario le falta una escena, el
# .get(img, img) deja pasar el nombre crudo del archivo y el documento sale con diecinueve escenas
# con nombre y una con «Triclobs_Kaptein_1123». Es exactamente lo que paso: cuando se sustituyo el
# par corrupto, DOS generadores quedaron nombrando el par retirado y sin el que lo reemplazo, y el
# Anexo 19 del informe salio en crudo. Nada lo delataba porque el .get no falla nunca.
print('\n=== 21. los diccionarios de nombres de escena cubren el corpus vivo ===')
try:
    from src.datasets import list_pairs as _lp
    _vivos = {Path(str(_v)).stem for _v, _i in _lp()}
except Exception as _e:                                          # pragma: no cover
    _vivos = set()
    ok(False, f'no se pudo obtener el corpus vivo ({type(_e).__name__}: {_e})')

if _vivos:
    for _arch in ('make_avances_report.py', 'make_reporte_optimos.py'):
        _p = RAIZ / 'experiments' / _arch
        if not _p.exists():
            continue
        _txt = _p.read_text(encoding='utf-8')
        # el diccionario es literal: se lo lee con ast, no con un regex de comillas
        _m = re.search(r'^ESCENA\s*=\s*(\{.*?^\})', _txt, re.S | re.M)
        if not _m:
            ok(False, f'{_arch}: no se encontro el diccionario ESCENA')
            continue
        try:
            _esc = ast.literal_eval(_m.group(1))
        except (SyntaxError, ValueError) as _e:
            ok(False, f'{_arch}: ESCENA no se pudo evaluar ({_e})')
            continue
        _sin_nombre = sorted(_vivos - set(_esc))
        _sobran = sorted(set(_esc) - _vivos)
        ok(not _sin_nombre, f'{_arch}: las {len(_vivos)} escenas del corpus tienen nombre legible'
                            + (f' — en crudo saldrian {_sin_nombre}' if _sin_nombre else ''))
        ok(not _sobran, f'{_arch}: ESCENA no nombra escenas ajenas al corpus'
                        + (f' — sobran {_sobran}' if _sobran else ''))

# --------------------------------------------- 20. docs/ raiz: solo entregables y estado vivo
# POR QUE. Los documentos de trabajo ya cumplidos se archivaron en docs/historial/ para que quien
# clone el repo no confunda un borrador viejo con el estado actual. Ese orden se deshace solo: basta
# que un generador siga escribiendo su salida en docs/ y el documento archivado reaparece arriba,
# con pinta de vigente, cada vez que alguien corre el script. Paso de verdad con
# make_reporte_optimos.py, que escribia docs/Resultados_Optimos_por_Imagen.pdf.
# El chequeo mira las dos caras: que en docs/ raiz solo esten los archivos declarados, y que ningun
# script escriba ahi un archivo que no sea uno de ellos.
print('\n--- 20. docs/ raiz tiene solo los entregables y el estado vivo')

_DOCS_RAIZ = {
    'Tesis_Borrador_V3.docx': 'entregable: el libro',
    'Tesis_Borrador_V3.pdf': 'entregable: el libro renderizado, para leer sin Word',
    'Avances_Tesis.pdf': 'entregable: el informe de avances',
    'Avances_Tesis_Tablas.xlsx': 'entregable: el libro de tablas',
    'Tesis_Defensa_Presentacion.pptx': 'entregable: el deck de defensa',
    'Tesis_Defensa_Presentacion.pdf': 'entregable: el deck renderizado',
    'ESTADO_Y_PENDIENTES.md': 'estado vivo del proyecto',
    'PROPUESTA_ENCUADRE.md': 'la redaccion propuesta para el encuadre del aporte, a decidir con el '
                            'director: no modifica ningun entregable',
    'PLAN_REESCRITURA_INFORME.md': 'que cambia en el informe al aclararse que el objetivo es uno y '
                                    'que el experimento que lo responde es el de M3FD',
    'AUDITORIA_LIBRO.md': 'las afirmaciones del libro marcadas al auditarlas contra los CSV, para la '
                          'revision manual: 987 revisadas, 910 respaldadas',
    'Auditoria_Bibliografia.md': 'vivo: lo reusa la auditoria de referencias',
    'Plan_Deck_Defensa.md': 'vivo: prescripciones del deck todavia sin aplicar',
}
_en_docs = sorted(p.name for p in (RAIZ / 'docs').iterdir() if p.is_file())
_sin_declarar = [n for n in _en_docs if n not in _DOCS_RAIZ]
_declarado_ausente = [n for n in _DOCS_RAIZ if n not in _en_docs]
ok(not _sin_declarar, f'los {len(_en_docs)} archivos de docs/ raiz estan declarados'
                      + (f' — sin declarar: {_sin_declarar}' if _sin_declarar else ''))
ok(not _declarado_ausente, f'estan los {len(_DOCS_RAIZ)} archivos que docs/ raiz debe tener'
                           + (f' — faltan: {_declarado_ausente}' if _declarado_ausente else ''))

# La otra cara: que ningun generador apunte su salida a docs/ raiz fuera de esa lista. Las rutas se
# escriben de dos formas en el repo —os.path.join(ROOT, "docs", "x.pdf") y el literal "docs/x.pdf",
# que en Windows a veces viene con barra invertida—, asi que la linea se normaliza a barra comun
# antes de mirarla y el patron no necesita ningun escape que se pueda perder al copiarlo.
_ESCRIBE = re.compile(
    r'''["']docs["']\s*,\s*["'](?P<a>[\w.\-]+\.(?:pdf|docx|xlsx|pptx|md|html))["']'''  # join("docs","x")
    r'''|docs/(?P<b>[\w.\-]+\.(?:pdf|docx|xlsx|pptx|html))["']''')                # el literal docs/algo.pdf
_fugas = {}
for _py in sorted(list((RAIZ / 'src').rglob('*.py')) + list((RAIZ / 'experiments').rglob('*.py'))):
    _txt = _py.read_text(encoding='utf-8', errors='replace')
    for _ln in _txt.splitlines():
        if _ln.lstrip().startswith('#') or 'historial' in _ln or '_local' in _ln:
            continue
        for _m in _ESCRIBE.finditer(_ln.replace('\\\\', '/').replace('\\', '/')):
            _n = _m.group('a') or _m.group('b')
            if _n not in _DOCS_RAIZ:
                _fugas.setdefault(_py.name, set()).add(_n)
ok(not _fugas, f'ningun script escribe en docs/ raiz fuera de los {len(_DOCS_RAIZ)} declarados'
               + (f' — {({k: sorted(v) for k, v in _fugas.items()})}' if _fugas else ''))

# --------------------------------------------- resumen
print(f'\n=== {len(fallos)} fallos · {len(avisos)} avisos ===')
for f in fallos:
    print(f'  FALLA {f}')
for a in avisos:
    print(f'  AVISO {a}')
sys.exit(1 if fallos else 0)
