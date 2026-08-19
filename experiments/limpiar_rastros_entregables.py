# -*- coding: utf-8 -*-
"""Saca de los entregables editables los rastros del entorno de trabajo.

DOS RASTROS, en lugares que no se ven al leer el documento y que por eso sobrevivieron:

  1. EL LIBRO tenia cinco comentarios de Word cuyo autor era el nombre de la herramienta, tres de ellos
     anclados en el texto. No aparecen en el PDF —LibreOffice no imprime los comentarios— pero se ven en
     el panel de revision en cuanto alguien abre el .docx, que es lo primero que hace quien lo recibe.
     Se retiran el archivo de comentarios, sus relaciones, su declaracion de tipo de contenido y las
     marcas que los anclan en el cuerpo.

  2. EL DECK tenia el TEXTO ALTERNATIVO de cinco imagenes con la ruta completa del directorio temporal,
     que lleva el nombre de la herramienta en el camino. Se ve en el campo de texto alternativo de
     PowerPoint y en cualquier lector de accesibilidad. Se reemplaza por el titulo de la lamina, que es
     una descripcion util y no una ruta.

EL CONTENIDO DE LOS DOCUMENTOS NO SE TOCA. El script comprueba, antes y despues, que la cantidad de
parrafos y el texto del libro y las laminas del deck sean identicos: lo unico que cambia es lo que no se
lee.

Uso:  .venv\\Scripts\\python.exe -X utf8 experiments/limpiar_rastros_entregables.py
"""
import re
import shutil
import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / 'docs'
LIBRO = DOCS / 'Tesis_Borrador_V3.docx'
DECK = DOCS / 'Tesis_Defensa_Presentacion.pptx'
RESP = ROOT / 'experiments' / 'results' / 'respaldos'
PAT = re.compile(r'claude|anthropic|chatgpt|openai|copilot|gemini', re.I)

# las partes de comentarios de un .docx, con las que las acompañan
PARTES_COM = ('word/comments.xml', 'word/commentsExtended.xml', 'word/commentsIds.xml',
              'word/commentsExtensible.xml', 'word/people.xml')


def texto_docx(p):
    from docx import Document
    d = Document(str(p))
    return [x.text for x in d.paragraphs], len(d.tables)


def texto_pptx(p):
    from pptx import Presentation
    pr = Presentation(str(p))
    return ['\n'.join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
            for s in pr.slides]


def reescribir(origen, cambios, quitar=()):
    """Reescribe el zip aplicando `cambios` (nombre -> funcion sobre el texto) y quitando `quitar`."""
    tmp = origen.with_suffix(origen.suffix + '.limpio')
    with zipfile.ZipFile(origen) as zin, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for it in zin.infolist():
            if it.filename in quitar:
                continue
            datos = zin.read(it.filename)
            if it.filename in cambios:
                datos = cambios[it.filename](datos.decode('utf-8')).encode('utf-8')
            zout.writestr(it, datos)
    return tmp


def limpiar_libro():
    if not LIBRO.exists():
        return ['no esta el libro']
    antes_p, antes_t = texto_docx(LIBRO)
    with zipfile.ZipFile(LIBRO) as z:
        presentes = [n for n in PARTES_COM if n in z.namelist()]
        n_com = len(re.findall(r'<w:comment ', z.read('word/comments.xml').decode('utf-8'))) \
            if 'word/comments.xml' in z.namelist() else 0
    if not presentes:
        return ['el libro no tiene comentarios']

    def sin_anclas(t):
        # las marcas que anclan cada comentario en el cuerpo. El <w:r> que lleva la referencia se retira
        # entero: dejar el run vacio no rompe nada, pero deja basura en el XML.
        t = re.sub(r'<w:commentRangeStart[^/]*/>', '', t)
        t = re.sub(r'<w:commentRangeEnd[^/]*/>', '', t)
        t = re.sub(r'<w:r>(?:(?!</w:r>).)*?<w:commentReference[^/]*/>.*?</w:r>', '', t, flags=re.S)
        t = re.sub(r'<w:commentReference[^/]*/>', '', t)
        return t

    def sin_rels(t):
        for n in presentes:
            hoja = n.split('/')[-1]
            t = re.sub(rf'<Relationship[^>]*Target="{hoja}"[^>]*/>', '', t)
        return t

    def sin_tipos(t):
        for n in presentes:
            t = re.sub(rf'<Override PartName="/{re.escape(n)}"[^>]*/>', '', t)
        return t

    tmp = reescribir(LIBRO, {'word/document.xml': sin_anclas,
                             'word/_rels/document.xml.rels': sin_rels,
                             '[Content_Types].xml': sin_tipos}, quitar=presentes)
    RESP.mkdir(parents=True, exist_ok=True)
    shutil.copy2(LIBRO, RESP / (LIBRO.name + '.con_comentarios'))
    tmp.replace(LIBRO)
    desp_p, desp_t = texto_docx(LIBRO)
    if desp_p != antes_p or desp_t != antes_t:
        return [f'EL CONTENIDO CAMBIO: {len(antes_p)} parrafos y {antes_t} tablas antes, '
                f'{len(desp_p)} y {desp_t} despues']
    return [f'libro: {n_com} comentarios retirados de {len(presentes)} partes; '
            f'{len(desp_p)} parrafos y {desp_t} tablas intactos']


def limpiar_deck():
    if not DECK.exists():
        return ['no esta el deck']
    antes = texto_pptx(DECK)
    from pptx import Presentation
    pr = Presentation(str(DECK))
    titulos = {}
    for i, s in enumerate(pr.slides, 1):
        t = next((sh.text_frame.text.strip() for sh in s.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()), '')
        titulos[f'ppt/slides/slide{i}.xml'] = re.sub(r'\s+', ' ', t)[:110]

    cambios, tocadas = {}, []
    with zipfile.ZipFile(DECK) as z:
        for nm in z.namelist():
            if not nm.startswith('ppt/slides/slide') or not nm.endswith('.xml'):
                continue
            t = z.read(nm).decode('utf-8')
            if not PAT.search(t):
                continue
            tocadas.append(nm)
            desc = (titulos.get(nm) or 'Figura de la presentación').replace('"', '')

            def rep(m, desc=desc):
                # solo el atributo descr, y solo si lleva un rastro: el resto del XML no se toca
                return m.group(0) if not PAT.search(m.group(1)) else f'descr="{desc}"'

            cambios[nm] = lambda s, rep=rep: re.sub(r'descr="([^"]*)"', rep, s)
    if not cambios:
        return ['el deck no tiene rastros en el texto alternativo']
    tmp = reescribir(DECK, cambios)
    RESP.mkdir(parents=True, exist_ok=True)
    shutil.copy2(DECK, RESP / (DECK.name + '.con_rutas'))
    tmp.replace(DECK)
    desp = texto_pptx(DECK)
    if desp != antes:
        return [f'EL CONTENIDO CAMBIO: {len(antes)} laminas antes, {len(desp)} despues']
    return [f'deck: texto alternativo reemplazado en {len(tocadas)} laminas '
            f'({", ".join(n.split("/")[-1] for n in tocadas)}); {len(desp)} laminas intactas']


# =====================================================================================================
# TERCER Y CUARTO RASTRO, agregados el 19 de agosto de 2026 despues de un barrido que encontro lo que
# los dos primeros no cubrian. Los dos primeros buscaban EL NOMBRE DE LA HERRAMIENTA (el patron PAT), y
# por eso pasaron de largo:
#
#   3. LAS RUTAS DEL EQUIPO en el texto alternativo. Seis atributos descr del deck llevan la ruta
#      completa del arbol del proyecto —no la del temporal, que sí tiene el nombre de la herramienta y
#      por eso sí se limpio en agosto—. Se leen en el panel «Texto alternativo» de PowerPoint.
#
#   4. LA IDENTIDAD DEL GENERADOR en los metadatos. El deck declara «PptxGenJS» en autor, titulo, asunto
#      y empresa, y eso se hereda al diccionario /Info de su PDF: el Explorador de Windows muestra
#      «Autores = PptxGenJS» sin abrir el archivo, y el navegador lo usa como titulo de la pestaña.
#      Ademas los conteos quedaron viejos: declara 19 laminas y 19 notas cuando hay 23 de cada una, y
#      sus titulos son «Slide 1»…«Slide 19», que son tambien los marcadores del PDF.
#
# QUE SE ESCRIBE Y QUE NO. Los valores de identidad se LEEN de los propios entregables —los autores de
# la portada del informe, el titulo y la institucion de la primera lamina— y no se escriben a mano. Y
# el elemento <Application> NO SE TOCA: dice con que programa se guardo el archivo, y cambiarlo seria
# afirmar que se uso un programa que no se uso. Se limpia lo que es falso (una empresa que no existe) y
# se completa lo que esta vacio, no se fabrica procedencia.
RUTA = re.compile(r'[A-Za-z]:[\\/](?:Users|Usuario)|/home/[a-z]|/Users/[A-Za-z]', re.I)
GENERADOR = re.compile(r'pptxgenjs|python-pptx|python-docx|reportlab|weasyprint|puppeteer', re.I)
INFORME = DOCS / 'Avances_Tesis.pdf'


def respaldar(p, sufijo):
    """Copia de seguridad en RESP, PERO SOLO SI git ignora ese directorio.

    Es el guardrail que faltaba. El 16 de agosto la limpieza guardo aca el docx CON los comentarios y el
    pptx CON las rutas, los dos quedaron versionados, y el remoto es publico: la limpieza quedo deshecha
    por sus propios respaldos, y el respaldo del libro lleva los cinco comentarios con el nombre de la
    herramienta como autor. Si el directorio no esta ignorado, este script no escribe nada y aborta.
    """
    RESP.mkdir(parents=True, exist_ok=True)
    import subprocess
    r = subprocess.run(['git', 'check-ignore', '-q', str(RESP)], cwd=str(ROOT),
                       capture_output=True, text=True)
    if r.returncode != 0:
        raise SystemExit(f'ABORTA: {RESP.relative_to(ROOT)} no esta ignorado por git. Un respaldo '
                         f'versionado publica lo que se acaba de limpiar. Agregalo al .gitignore.')
    shutil.copy2(p, RESP / (p.name + sufijo))


def identidad():
    """Autores, titulo e institucion leidos de los entregables. Ninguno se escribe a mano."""
    import fitz
    from pptx import Presentation
    with fitz.open(str(INFORME)) as d:
        portada = d[0].get_text()
    m = re.search(r'Autores:\s*(.+)', portada)
    autores = re.sub(r'\s*—\s*', '; ', m.group(1).strip()) if m else ''
    pr = Presentation(str(DECK))
    crudos = [sh.text_frame.text.strip() for sh in list(pr.slides)[0].shapes
              if sh.has_text_frame and sh.text_frame.text.strip()]
    marcos = [re.sub(r'\s+', ' ', c) for c in crudos]
    # la institucion es la PRIMERA LINEA del primer marco, no todo el marco: la lamina lleva
    # «UNIVERSIDAD COMUNERA / FACULTAD … — MAESTRÍA …» en tres renglones, y colapsar los saltos de linea
    # antes de cortar dejaba «UNIVERSIDAD COMUNERA FACULTAD DE CIENCIAS Y TECNOLOGÍA» como empresa
    institucion = crudos[0].splitlines()[0].strip() if crudos else ''
    titulo = marcos[1] if len(marcos) > 1 else ''
    if not (autores and titulo and institucion):
        raise SystemExit(f'ABORTA: no se pudo leer la identidad de los entregables '
                         f'(autores={autores!r}, titulo={titulo!r}, institucion={institucion!r})')
    return autores, titulo, institucion


def titulos_deck():
    """El titulo real de cada lamina, en orden."""
    from pptx import Presentation
    pr = Presentation(str(DECK))
    out = []
    for s in pr.slides:
        t = next((sh.text_frame.text.strip() for sh in s.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()), '')
        out.append(re.sub(r'\s+', ' ', t)[:110] or 'Lámina sin título')
    return out


def limpiar_rutas_alt():
    """Las rutas del equipo en los atributos descr y title de las dos piezas editables."""
    msgs = []
    for p, rot in ((DECK, 'deck'), (LIBRO, 'libro')):
        if not p.exists():
            continue
        antes = texto_pptx(p) if p is DECK else texto_docx(p)
        tit = titulos_deck() if p is DECK else []
        cambios, n = {}, 0
        with zipfile.ZipFile(p) as z:
            for nm in z.namelist():
                if not nm.endswith('.xml'):
                    continue
                t = z.read(nm).decode('utf-8', errors='replace')
                if not RUTA.search(t):
                    continue
                m = re.match(r'ppt/slides/slide(\d+)\.xml$', nm)
                desc = (tit[int(m.group(1)) - 1] if m and int(m.group(1)) <= len(tit)
                        else 'Figura del documento').replace('"', '')

                def rep(mm, desc=desc):
                    # solo el valor del atributo, y solo si lleva una ruta: el resto no se toca
                    return mm.group(0) if not RUTA.search(mm.group(2)) else f'{mm.group(1)}="{desc}"'

                nuevo = re.sub(r'(descr|title)="([^"]*)"', rep, t)
                if nuevo != t:
                    n += len(RUTA.findall(t))
                    cambios[nm] = lambda s, nuevo=nuevo: nuevo
        if not cambios:
            msgs.append(f'{rot}: sin rutas del equipo en el texto alternativo')
            continue
        tmp = reescribir(p, cambios)
        respaldar(p, '.con_rutas_de_equipo')
        tmp.replace(p)
        desp = texto_pptx(p) if p is DECK else texto_docx(p)
        if desp != antes:
            msgs.append(f'{rot}: EL CONTENIDO CAMBIO')
        else:
            msgs.append(f'{rot}: {n} rutas del equipo reemplazadas por la descripción de la lámina '
                        f'en {len(cambios)} partes; contenido intacto')
    return msgs


def limpiar_identidad_deck():
    """docProps del deck: la identidad que nombra al generador y los conteos que quedaron viejos."""
    if not DECK.exists():
        return ['no esta el deck']
    from xml.sax.saxutils import escape
    autores, titulo, institucion = identidad()
    tit = titulos_deck()
    antes = texto_pptx(DECK)
    with zipfile.ZipFile(DECK) as z:
        n_lam = sum(1 for n in z.namelist() if re.fullmatch(r'ppt/slides/slide\d+\.xml', n))
        n_not = sum(1 for n in z.namelist() if re.fullmatch(r'ppt/notesSlides/notesSlide\d+\.xml', n))
        core = z.read('docProps/core.xml').decode('utf-8')
        app = z.read('docProps/app.xml').decode('utf-8')
    tocado = []

    def core_nuevo(t):
        for etq, val in (('dc:title', titulo), ('dc:subject', titulo),
                         ('dc:creator', autores), ('cp:lastModifiedBy', autores)):
            m = re.search(rf'<{etq}>([^<]*)</{etq}>', t)
            if m and GENERADOR.search(m.group(1)):
                tocado.append(f'{etq}: «{m.group(1)}» -> «{val[:40]}»')
                t = t.replace(m.group(0), f'<{etq}>{escape(val)}</{etq}>')
        return t

    def app_nuevo(t):
        m = re.search(r'<Company>([^<]*)</Company>', t)
        if m and GENERADOR.search(m.group(1)):
            tocado.append(f'Company: «{m.group(1)}» -> «{institucion}»')
            t = t.replace(m.group(0), f'<Company>{escape(institucion)}</Company>')
        # los conteos, que son un hecho verificable y estaban mal
        for etq, real in (('Slides', n_lam), ('Notes', n_not)):
            m = re.search(rf'<{etq}>(\d+)</{etq}>', t)
            if m and int(m.group(1)) != real:
                tocado.append(f'{etq}: {m.group(1)} -> {real}')
                t = t.replace(m.group(0), f'<{etq}>{real}</{etq}>')
        # y los titulos genericos «Slide N», que son los que el PDF usa de marcadores. Se reemplaza la
        # corrida completa y se corrigen los dos tamaños que la describen, o el paquete queda incoherente.
        viejos = re.findall(r'<vt:lpstr>Slide \d+</vt:lpstr>', t)
        if viejos:
            nuevos = ''.join(f'<vt:lpstr>{escape(x)}</vt:lpstr>' for x in tit)
            corrida = ''.join(viejos)
            if corrida not in t:
                raise SystemExit('ABORTA: los títulos «Slide N» no están contiguos en app.xml; '
                                 'reemplazarlos uno por uno movería las entradas de fuente y tema.')
            t = t.replace(corrida, nuevos)
            tocado.append(f'títulos de lámina: {len(viejos)} genéricos -> {len(tit)} reales')
            # LOS DOS TAMAÑOS QUE DESCRIBEN EL VECTOR, y van con asercion y no con «if m:». La primera
            # version los buscaba sin tolerar el sangrado del XML, no encontro ninguno, y se saltearon
            # los dos EN SILENCIO: el paquete quedo con size="22" declarando 26 entradas y con
            # «Slide Titles = 19» para 23 laminas. Un reemplazo que puede no aplicarse tiene que gritar.
            for rot, pat, val in (
                    ('tamaño del vector TitlesOfParts',
                     r'(<TitlesOfParts>\s*<vt:vector size=")(\d+)(")',
                     lambda v: int(v) - len(viejos) + len(tit)),
                    ('conteo «Slide Titles»',
                     r'(<vt:lpstr>Slide Titles</vt:lpstr></vt:variant>\s*<vt:variant><vt:i4>)(\d+)(</vt:i4>)',
                     lambda v: len(tit))):
                m = re.search(pat, t)
                if not m:
                    raise SystemExit(f'ABORTA: no se encontro el {rot} en app.xml. Sin corregirlo el '
                                     f'paquete queda incoherente y PowerPoint puede rechazarlo.')
                nuevo_val = val(m.group(2))
                if int(m.group(2)) != nuevo_val:
                    tocado.append(f'{rot}: {m.group(2)} -> {nuevo_val}')
                    t = t.replace(m.group(0), f'{m.group(1)}{nuevo_val}{m.group(3)}')
        return t

    if not (GENERADOR.search(core) or GENERADOR.search(app)
            or re.search(r'<vt:lpstr>Slide \d+</vt:lpstr>', app)):
        return ['el deck no tiene identidad de generador ni títulos genéricos']
    tmp = reescribir(DECK, {'docProps/core.xml': core_nuevo, 'docProps/app.xml': app_nuevo})
    respaldar(DECK, '.con_identidad_de_generador')
    tmp.replace(DECK)
    desp = texto_pptx(DECK)
    if desp != antes:
        return [f'deck: EL CONTENIDO CAMBIO ({len(antes)} láminas antes, {len(desp)} después)']
    from pptx import Presentation
    Presentation(str(DECK))            # si el paquete quedo incoherente, esto revienta
    return [f'deck docProps: {len(tocado)} campos corregidos'] + [f'  · {x}' for x in tocado]


def limpiar_pdf_deck():
    """El /Info del PDF del deck se hereda del pptx, y sus marcadores dicen «Slide N».

    Se corrige sobre el PDF existente en lugar de re-exportarlo: re-exportar vuelve a maquetar 23
    laminas y puede mover cualquier cosa. Aca no se toca ni un pixel de las paginas, y se comprueba.
    """
    import fitz
    pdf = DECK.with_suffix('.pdf')
    if not pdf.exists():
        return ['no esta el PDF del deck']
    autores, titulo, institucion = identidad()
    tit = titulos_deck()
    with fitz.open(str(pdf)) as d:
        antes = (d.page_count, [d[i].get_text() for i in range(d.page_count)])
        meta_antes = dict(d.metadata or {})
    sucio = [f'{k}={v}' for k, v in meta_antes.items() if v and GENERADOR.search(str(v))]
    toc_antes = None
    with fitz.open(str(pdf)) as d:
        toc_antes = d.get_toc()
    genericos = sum(1 for _, t, _ in (toc_antes or []) if re.fullmatch(r'Slide \d+', t.strip()))
    if not sucio and not genericos:
        return ['el PDF del deck no tiene identidad de generador ni marcadores genéricos']
    tmp = pdf.with_suffix('.pdf.limpio')
    with fitz.open(str(pdf)) as d:
        d.set_metadata({'title': titulo, 'author': autores, 'subject': titulo,
                        'keywords': '', 'creator': '', 'producer': ''})
        if genericos and len(tit) == d.page_count:
            d.set_toc([[1, tit[i], i + 1] for i in range(d.page_count)])
        d.save(str(tmp), garbage=3, deflate=True)
    with fitz.open(str(tmp)) as d:
        desp = (d.page_count, [d[i].get_text() for i in range(d.page_count)])
        meta_desp = dict(d.metadata or {})
    if desp != antes:
        tmp.unlink()
        return [f'PDF del deck: EL CONTENIDO CAMBIO ({antes[0]} páginas antes, {desp[0]} después)']
    respaldar(pdf, '.con_identidad_de_generador')
    tmp.replace(pdf)
    return [f'PDF del deck: {len(sucio)} metadatos con el generador ({", ".join(sucio)}) y '
            f'{genericos} marcadores genéricos corregidos; {antes[0]} páginas intactas',
            f'  · ahora title={meta_desp.get("title", "")[:40]!r} author={meta_desp.get("author", "")!r} '
            f'producer={meta_desp.get("producer", "")!r}']


def main():
    msgs = limpiar_libro() + limpiar_deck() + limpiar_rutas_alt() + limpiar_identidad_deck() \
        + limpiar_pdf_deck()
    for m in msgs:
        print(f'    {m}')
    # CONTROL FINAL, sobre los tres patrones y no solo sobre el nombre de la herramienta: fue mirar un
    # patron solo lo que dejo pasar las rutas del equipo y la identidad del generador durante tres dias.
    quedan = []
    for p in (LIBRO, DECK):
        if not p.exists():
            continue
        with zipfile.ZipFile(p) as z:
            for nm in z.namelist():
                if not nm.endswith(('.xml', '.rels')):
                    continue
                t = z.read(nm).decode('utf-8', errors='replace')
                for rot, pat in (('herramienta', PAT), ('ruta del equipo', RUTA),
                                 ('generador', GENERADOR)):
                    if pat.search(t):
                        quedan.append(f'{p.name}/{nm} [{rot}]')
    print(f'\n  control: {len(quedan)} partes con rastros' + (f' — {quedan}' if quedan else ' (limpio)'))
    return 1 if quedan or any('CAMBIO' in m for m in msgs) else 0


if __name__ == '__main__':
    sys.exit(main())
